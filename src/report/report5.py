from functools import reduce

from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE, XL_LABEL_POSITION
from pptx.util import Cm  # Inches
from pandas import np
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.chart import XL_LEGEND_POSITION

from pptx.dml.color import RGBColor
import pandas as pd

import report_common

if __name__ == '__main__':

    df_vw = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 0)
    df_mkt = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 1)
    prs = Presentation('c:/auto-report/cover.pptx')
    title_only_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    print("df_vw=" + str(df_vw.shape[0]))
    print("df_mkt=" + str(df_mkt.shape[0]))
    view_rules = 'OEM'  # or OEM
    all_mkt = False  # true分母是否为全市场，false分母为fuel_type细分市场
    oem = df_vw.groupby('OEM').size().reset_index()["OEM"]
    brand = df_vw.groupby('Brand').size().reset_index()["Brand"]
    fuel_type = df_vw.groupby('Fuel_Type').size().reset_index()["Fuel_Type"]
    fuel_type_group = df_vw.groupby('Fuel_Type_Group').size().reset_index()["Fuel_Type_Group"]

    print(oem)
    print(brand)
    print(fuel_type)
    print(fuel_type_group)
    # 此处应该对 oem brand fuel_type三个数组进行过滤，把用户不需要的删除掉，默认是全选，需要通过读取web端的配置文件
    # oem = ['FAW-VW', 'JAC-VW', 'JV TBD', 'SAIC-VW']
    # brand = ['Audi', 'Cupra', 'Jetta', 'Sihao', 'Skoda', 'VW']
    # fuel_type = ['ICE', 'BEV', 'PHEV']
    oem_filter = oem
    brand_filter = brand
    fuel_type_filter = fuel_type  # 根据私有fitler决定哪个fuel_type需要保留

    PR_Status_local = 'PR67.OP'  # 本轮
    PR_Status_previous = 'PR67.SP'  # 对比轮也叫上一轮
    PR_Status = []
    PR_Status.append(PR_Status_local)
    PR_Status.append(PR_Status_previous)
    start_year = 2018
    year_span = 9
    end_year = start_year + year_span
    # oem = ['FAW-VW', 'JAC-VW', 'JV TBD', 'SAIC-VW']
    # brand = ['Audi', 'Cupra', 'Jetta', 'Sihao', 'Skoda', 'VW']
    # 根据所有filter过滤大众数据
    df_vw_filter = \
        df_vw[(df_vw['PR_Status'].isin(PR_Status)) & (df_vw['YEAR'] >= start_year) & (df_vw['YEAR'] <= end_year) \
              & (df_vw['OEM'].isin(oem)) & (df_vw['Brand'].isin(brand)) & (df_vw['Fuel_Type'].isin(fuel_type))]
    df_mkt_filter = \
        df_mkt[(df_mkt['Status'].isin(PR_Status)) & (df_mkt['Year'] >= start_year) & (df_mkt['Year'] <= end_year)]
    print("df_vw_filter=" + str(df_vw_filter.shape[0]))
    print("df_mkt_filter=" + str(df_mkt_filter.shape[0]))

    # 获得需要显示的年份数组，在图标中最为重要y轴坐标
    data_years = df_vw_filter.groupby(['YEAR']).size().reset_index().sort_values(['YEAR'], ascending=[True])['YEAR']
    # print(data_years)

    # 计算本轮和对比轮的总量
    df_vw_status = df_vw_filter.groupby(['PR_Status']).agg({'Volume': np.sum}).reset_index()
    df_vw_status_fuel_year = df_vw_filter.groupby(['PR_Status', 'Fuel_Type', 'YEAR']).agg({'Volume': np.sum}).reset_index()
    PR_Status_local_vol = df_vw_status[df_vw_status['PR_Status'] == PR_Status_local].reset_index().loc[0, 'Volume']
    PR_Status_previous_vol = df_vw_status[df_vw_status['PR_Status'] == PR_Status_previous].reset_index() \
        .loc[0, 'Volume']
    print(PR_Status_local_vol)
    print(PR_Status_previous_vol)

    # 销售市场数据分组聚合，并按Fuel_type和year分组计算market effect rate(ice bev phev)-------------------------
    df_mkt_volume_local = \
        df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_local)].groupby(['Fuel_type', 'Year']) \
        .agg({'Volume': np.sum}).reset_index()
    print(df_mkt_volume_local)
    df_mkt_volume_previous = \
        df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_previous)].groupby(['Fuel_type', 'Year']) \
        .agg({'Volume': np.sum}).reset_index()
    print(df_mkt_volume_previous)

    year_mkt_rate_dict = {}
    year_mkt_effect_dict = {}
    year_mkt_effect_list = []
    total_mkt_effect_volume = 0
    total_mkt_effect_percent = 0
    for year in data_years:
        fuel_mkt_rate_dict = {}
        fuel_mkt_effect_dict = {}
        fuel_mkt_effect_sum = 0
        for fuel in fuel_type_filter:
            # print('maoyadong+' + str(year) + '+' + str(fuel))
            # print(df_mkt_volume_local[(df_mkt_volume_local['Year'] == year) & (df_mkt_volume_local['Fuel_type'] == fuel)].reset_index())
            mkt_volume_local = \
                df_mkt_volume_local[(df_mkt_volume_local['Year'] == year) & (df_mkt_volume_local['Fuel_type'] == fuel)] \
                .reset_index().loc[0, 'Volume']
            mkt_volume_previous = \
                df_mkt_volume_previous[(df_mkt_volume_local['Year'] == year) & (df_mkt_volume_local['Fuel_type'] == fuel)] \
                .reset_index().loc[0, 'Volume']
            mkt_rate = mkt_volume_local / mkt_volume_previous - 1
            vw_volume_previous = df_vw_status_fuel_year[ \
                (df_vw_status_fuel_year['PR_Status'] == PR_Status_previous) & \
                (df_vw_status_fuel_year['YEAR'] == year) & \
                (df_vw_status_fuel_year['Fuel_Type'] == fuel)] \
                .reset_index().loc[0, 'Volume']
            mkt_effect_volume = vw_volume_previous * mkt_rate
            fuel_mkt_rate_dict[fuel] = mkt_rate
            fuel_mkt_effect_dict[fuel] = mkt_effect_volume
            fuel_mkt_effect_sum = fuel_mkt_effect_sum + mkt_effect_volume
        year_mkt_rate_dict[year] = fuel_mkt_rate_dict
        year_mkt_effect_dict[year] = fuel_mkt_effect_dict
        year_mkt_effect_list.append(fuel_mkt_effect_sum)
        total_mkt_effect_volume = total_mkt_effect_volume + fuel_mkt_effect_sum
    total_mkt_effect_percent = total_mkt_effect_percent / PR_Status_previous_vol
    print('mkt_rate_dict=' + str(year_mkt_rate_dict))
    print('year_mkt_effect_dict=' + str(year_mkt_effect_dict))
    print('year_mkt_effect_list' + str(year_mkt_effect_list))
    print('total_mkt_effect_volume=' + str(total_mkt_effect_volume))


    # 按大众本轮的销量排序OEM和Brand
    categories_order = \
        df_vw_filter[df_vw_filter['PR_Status'].isin(PR_Status)].groupby(view_rules).agg({'Volume': np.sum}) \
        .reset_index().sort_values(['Volume'], ascending=[False]).reset_index()[view_rules]
    print(categories_order)


    # 计算每year每个Brand或OEM 每个Fuel_Type的 market effect volume(ice bev phev)，然后在按brand或OEM相加-------------------------
    # 计算每year每个Brand或OEM的change,每年本轮-每年对比轮-每年Market effect volume
    vw_vol_local = df_vw_filter[df_vw_filter['PR_Status'] == PR_Status_local] \
        .groupby([view_rules, 'Fuel_Type', 'YEAR']).agg({'Volume': np.sum}).reset_index()
    vw_vol_previous = df_vw_filter[df_vw_filter['PR_Status'] == PR_Status_previous] \
        .groupby([view_rules, 'Fuel_Type', 'YEAR']).agg({'Volume': np.sum}).reset_index()

    category_change_dict = {}
    category_change_list = []
    category_change_precent_list = []

    for category in categories_order:
        year_volume_change_sum = 0
        year_volume_change_list = []
        for year in data_years:
            fuel_change_sum = 0
            for fuel in fuel_type_filter:
                fuel_vol_local = 0
                fuel_vol_previous = 0
                if not vw_vol_local[ \
                        (vw_vol_local[view_rules] == category) & (vw_vol_local['Fuel_Type'] == fuel) & \
                        (vw_vol_local['YEAR'] == year)].empty:
                    fuel_vol_local = vw_vol_local[ \
                        (vw_vol_local[view_rules] == category) & (vw_vol_local['Fuel_Type'] == fuel) & \
                        (vw_vol_local['YEAR'] == year)].agg({'Volume': np.sum}).reset_index().iloc[0, 1]
                if not vw_vol_previous[ \
                        (vw_vol_previous[view_rules] == category) & (vw_vol_previous['Fuel_Type'] == fuel) & \
                        (vw_vol_previous['YEAR'] == year)].empty:
                    fuel_vol_previous = vw_vol_previous[
                        (vw_vol_previous[view_rules] == category) & (vw_vol_previous['Fuel_Type'] == fuel) & \
                        (vw_vol_previous['YEAR'] == year)].agg({'Volume': np.sum}).reset_index().iloc[0, 1]
                fuel_change = fuel_vol_local - fuel_vol_previous - fuel_vol_previous * year_mkt_rate_dict[year][fuel]
                fuel_change_sum = fuel_change_sum + fuel_change
                # print(category+str(year)+fuel+'#'+str(fuel_change)+'###'+str(fuel_vol_local)+'-'+str(fuel_vol_previous)+'-'+str(fuel_vol_previous)+'*'+str(year_mkt_rate_dict[year][fuel]))
            year_volume_change_list.append(fuel_change_sum)
            year_volume_change_sum = year_volume_change_sum + fuel_change_sum
        category_change_dict[category] = year_volume_change_list
        category_change_list.append(year_volume_change_sum)
        category_change_precent_list.append(year_volume_change_sum/PR_Status_previous_vol)
    print('category_change_dict=' + str(category_change_dict))
    print('category_change_list=' + str(category_change_list))
    print('category_change_precent_list=' + str(category_change_precent_list))

    # print("对比两种方法计算的total effect 和 category_change 是否相等")
    # category_change_sum = []
    # for idx, val in year_category_change_dict.items():
    #     category_change_sum.append(reduce(lambda x, y: x + y, val))
    #
    # year_mkt_total_effect_sum = reduce(lambda x, y: x + y, year_mkt_total_effect_list)

    #设置字体大小
    font_size = Pt(10)

    # 设置堆积图y轴坐标和
    data_categories = []
    # categories = [cat + ' change' for cat in categories_order]
    data_categories.extend(categories_order)
    data_categories.insert(0, PR_Status_previous)
    data_categories.insert(1, 'mkt effect')
    data_categories.append(PR_Status_local)
    print(data_categories)

    chart_data_stack = CategoryChartData()
    chart_data_stack.categories = data_categories

    # 设置柱状堆积图的series数据---------------------------------
    series_Volumes1 = []
    series_Volumes2 = []
    gap_sum = total_mkt_effect_volume

    for effect in category_change_list:
        gap_sum += effect
        gap = PR_Status_previous_vol + gap_sum
        series_Volumes1.append(gap)
        series_Volumes2.append(abs(effect))

    series_Volumes1.insert(0, 0)
    series_Volumes1.insert(1, PR_Status_previous_vol + total_mkt_effect_volume)
    series_Volumes1.append(0)

    series_Volumes2.insert(0, PR_Status_previous_vol)
    series_Volumes2.insert(1, abs(total_mkt_effect_volume))
    series_Volumes2.append(PR_Status_local_vol)

    series_Volumes2 = [vol / 1000 for vol in series_Volumes2]
    series_Volumes1 = [vol / 1000 for vol in series_Volumes1]
    chart_data_stack.add_series('series1', series_Volumes1)
    chart_data_stack.add_series('series2', series_Volumes2)
    print(series_Volumes1)
    print(series_Volumes2)

    x, y, cx, cy = Cm(1), Cm(3.5), Cm(24), Cm(6)
    graphic_frame_stack = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data_stack
    )

    chart_stack = graphic_frame_stack.chart
    # chart_stack.has_title = True
    # chart_stack.chart_title.has_text_frame = True
    # chart_stack.chart_title.text_frame.text = "maoyadong"
    # chart_stack.chart_title.text_frame.paragraphs[0].font.size = Pt(10)

    # chart_stack.has_legend = True
    # chart_stack.legend.position = XL_LEGEND_POSITION.LEFT  # XL_LEGEND_POSITION.CORNER
    # chart_stack.legend.include_in_layout = False
    # chart_stack.legend.font.size = Pt(10)

    chart_stack.series[0].data_labels.show_value = False
    chart_stack.series[0].format.fill.background()
    chart_stack.series[1].data_labels.show_value = True
    chart_stack.series[1].data_labels.number_format = '0'
    chart_stack.series[1].data_labels.font.size = font_size
    # chart_stack.series[1].data_labels.position = XL_LABEL_POSITION.ABOVE
    # for p in chart_stack.series[1].points:
    #     p.data_label.font.size = Pt(20)
    #     p.data_label.number_format = '0'

    for idx, category in enumerate(category_change_list):
        index = idx + 2
        series_point = chart_stack.series[1].points[index]
        series_point.data_label.has_text_frame = True
        paragraphs1 = series_point.data_label.text_frame.paragraphs[0]
        paragraphs2 = series_point.data_label.text_frame.add_paragraph()
        run1 = paragraphs1.add_run()
        run1.text = format(category_change_list[idx] / 1000, '.0f')
        run1.font.size = font_size
        run2 = paragraphs2.add_run()
        run2.text = format(category_change_precent_list[idx], '.1%')
        run2.font.size = font_size

    series_point = chart_stack.series[1].points[1]
    series_point.data_label.has_text_frame = True
    paragraphs1 = series_point.data_label.text_frame.paragraphs[0]
    paragraphs2 = series_point.data_label.text_frame.add_paragraph()
    run1 = paragraphs1.add_run()
    run1.text = format(total_mkt_effect_volume / 1000, '.0f')
    run1.font.size = font_size
    run2 = paragraphs2.add_run()
    run2.text = format(total_mkt_effect_percent, '.1%')
    run2.font.size = font_size
    # chart_stack.series[1].points[2].data_label.text_frame.paragraphs[1].text = 'paragraphs2'

    #设置中间柱子的颜色 负数为红 正数为绿
    for idx, series_point in enumerate(chart_stack.series[1].points):
        if 0 < idx < len(chart_stack.series[1].points) - 1:
            series_point.format.fill.patterned() #设置为可更改填充颜色的
            series_point.format.fill.solid() #设置为纯色填充
            if float(series_point.data_label.text_frame.paragraphs[0].runs[0].text) > 0:
                series_point.format.fill.fore_color.rgb = RGBColor(34, 139, 34) #绿色
            else:
                series_point.format.fill.fore_color.rgb = RGBColor(220, 20, 60) #红色


    value_axis_stack = chart_stack.value_axis
    value_axis_stack.has_major_gridlines = False
    value_axis_stack.major_tick_mark = XL_TICK_MARK.NONE
    value_axis_stack.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis_stack.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    value_axis_stack.visible = False

    category_axis_stack = chart_stack.category_axis
    category_axis_stack.has_major_gridlines = False
    category_axis_stack.major_tick_mark = XL_TICK_MARK.NONE
    category_axis_stack.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    category_axis_stack.tick_labels.font.size = Pt(8)
    category_axis_stack.format.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
    category_axis_stack.visible = True

    # value_axis_stack.has_title = True
    # value_axis_stack.axis_title.has_text_frame = True
    # value_axis_stack.axis_title.text_frame.text = "False positive"
    # value_axis_stack.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

    # 开始创建表格
    rows = len(categories_order) + 2

    cols = len(data_years) + 1
    table_width = 24
    table_height = 3
    top = Cm(9.5)
    left = Cm(1.5)  # Inches(2.0)
    width = Cm(table_width)  # Inches(6.0)
    height = Cm(table_height)  # Inches(0.8)

    # 添加表格到幻灯片 --------------------
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # 给data_year加E
    data_years_E = [str(year) + 'E' if year > start_year else str(year) for year in data_years]

    # 设置单元格宽度
    columns_width = table_width / cols - 0.1
    for i in range(cols):
        if i == 0:
            table.columns[i].width = Cm(columns_width + 0.6)
        else:
            table.columns[i].width = Cm(columns_width)  # Inches(2.0)

    row_height = table_height / rows
    for i in range(rows):
        table.rows[i].height = Cm(row_height)  # Inches(2.0)

    # 设置标题行
    for i in range(cols):
        if i == 0:
            table.cell(0, i).text = "In'000 units"
        else:
            table.cell(0, i).text = str(data_years_E[i - 1])

    # 设置第二行内容
    for i in range(cols):
        if i == 0:
            table.cell(1, i).text = "Market Effect"
        else:
            table.cell(1, i).text = format(year_mkt_effect_list[i - 1] / 1000, '.0f')

        # 填充表格第三行到最后一行数据
        for row_idx, category in enumerate(categories_order):
            table.cell(row_idx + 2, 0).text = category
            for col_idx, vol in enumerate(category_change_dict[category]):
                table.cell(row_idx + 2, col_idx + 1).text = format(vol / 1000, '.0f')


    # 调整table每个cell的字体
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(table):
        if cell.text.strip() == '':
            cell.text = r'/'
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = font_size

# 添加ppt顶部的椭圆形shapes
    shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(12.5), Cm(2.3), Cm(2), Cm(1.5))
    report_common.set_shape_oval_format(shape, font_size, PR_Status_local_vol - PR_Status_previous_vol, \
                                        (PR_Status_local_vol - PR_Status_previous_vol) / PR_Status_previous_vol)
    # shape.text_frame.word_wrap = False
    # shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    # paragraphs1 = shape.text_frame.paragraphs[0]
    # paragraphs2 = shape.text_frame.add_paragraph()
    # run1 = paragraphs1.add_run()
    # run1.text = format((PR_Status_local_vol - PR_Status_previous_vol) / 1000, '.0f')
    # run1.font.size = font_size
    # run1.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    # run2 = paragraphs2.add_run()
    # run2.text = format((PR_Status_local_vol - PR_Status_previous_vol) / PR_Status_previous_vol, '.1%')
    # run2.font.size = font_size
    # run2.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    # shape.fill.background()
    # shape.line.width = Pt(0.2)

    # 开始添加单位注释文本框
    # 在指定位置添加文本框
    textbox = shapes.add_textbox(Cm(1), Cm(2), Cm(2), Cm(1))
    tf = textbox.text_frame

    # 在文本框中写入文字
    para = tf.add_paragraph()  # 新增段落
    para.text = "Volume '000units"  # 向段落写入文字
    para.line_spacing = 1.5  # 1.5 倍的行距
    para.font.size = Pt(6)


    # shape2 = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(13), Cm(1), Cm(1.5), Cm(1.5))
    # shape2.text = '2'
    # shape2.fill.background()
    # shape3 = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(26), Cm(4), Cm(1), Cm(1))
    # shape3.text = '3'
    # shape3.fill.background()
    #
    # connector = shapes.add_connector(
    #     MSO_CONNECTOR.ELBOW, 1, 1, 1, 1
    # )
    # connector2 = shapes.add_connector(
    #     MSO_CONNECTOR.ELBOW, 1, 1, 1, 1
    # )
    # connector.begin_connect(shape, 0)
    # connector.end_connect(shape2, 1)
    # connector2.begin_connect(shape2, 3)
    # connector2.end_connect(shape3, 0)

    prs.save('c:/auto-report/template_tmp5.pptx')
    print("maoyadong")
