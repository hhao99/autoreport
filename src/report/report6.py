import math
from collections import Counter, OrderedDict
from functools import reduce

from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE, XL_LABEL_POSITION
from pptx.util import Cm  # Inches
from pandas import np
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.chart import XL_LEGEND_POSITION

from pptx.dml.color import RGBColor
import pandas as pd
import report_common

if __name__ == '__main__':

    df_vw = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 2)
    df_mkt = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 3)
    df_sale = pd.read_excel(r'c:/auto-report/sale_month.xlsx', 0)
    prs = Presentation('c:/auto-report/cover.pptx')
    title_only_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    print("df_vw=" + str(df_vw.shape[0]))
    print("df_mkt=" + str(df_mkt.shape[0]))
    view_rules = 'OEM'  # or OEM
    all_mkt = False  # true分母是否为全市场，false分母为fuel_type细分市场
    oem = df_vw.groupby('OEM').size().reset_index()["OEM"]
    brand = df_vw.groupby('Brand').size().reset_index()["Brand"]
    fuel_type = df_vw.groupby('Fuel_Type').size().reset_index()["Fuel_Type"]
    fuel_type_group = df_vw.groupby('Fuel_Type_Group').size().reset_index()["Fuel_Type_Group"]
    build_type = df_vw.groupby('Build_Type').size().reset_index()["Build_Type"]

    print(oem)
    print(brand)
    print(fuel_type)
    print(fuel_type_group)
    print(build_type)

    # 此处应该对 oem brand fuel_type build_type几个数组进行过滤，把用户不需要的删除掉，默认是全选，需要通过读取web端的配置文件
    # oem = ['FAW-VW', 'JAC-VW', 'JV TBD', 'SAIC-VW']
    # brand = ['Audi', 'Cupra', 'Jetta', 'Sihao', 'Skoda', 'VW']
    # fuel_type = ['ICE', 'BEV', 'PHEV']
    oem_filter = oem
    brand_filter = brand
    fuel_type_filter = fuel_type  # 根据私有fitler决定哪个fuel_type需要保留
    build_type_filter = build_type

    PR_Status_local = 'PR66.OP'  # 本轮
    PR_Status_previous = 'PR66.SP'  # 对比轮,也叫上一轮
    PR_Status = []
    PR_Status.append(PR_Status_local)
    PR_Status.append(PR_Status_previous)
    start_year = 2018
    year_span = 9
    end_year = start_year + year_span
    # oem = ['FAW-VW', 'JAC-VW', 'JV TBD', 'SAIC-VW']
    # brand = ['Audi', 'Cupra', 'Jetta', 'Sihao', 'Skoda', 'VW']
    # 根据所有filter过滤大众数据
    df_vw_filter = \
        df_vw[(df_vw['PR_Status'].isin(PR_Status)) & (df_vw['YEAR'] >= start_year) & (df_vw['YEAR'] <= end_year) \
              & (df_vw['OEM'].isin(oem_filter)) & (df_vw['Brand'].isin(brand_filter)) \
              & (df_vw['Fuel_Type'].isin(fuel_type_filter)) & (df_vw['Build_Type'].isin(build_type_filter))]
    df_mkt_filter = \
        df_mkt[(df_mkt['Status'].isin(PR_Status)) & (df_mkt['Year'] >= start_year) & (df_mkt['Year'] <= end_year)]
    print("df_vw_filter=" + str(df_vw_filter.shape[0]))
    print("df_mkt_filter=" + str(df_mkt_filter.shape[0]))

    # 分别获得大众两轮的全字段数据
    df_vw_local = df_vw_filter[df_vw_filter['PR_Status'] == PR_Status_local].reset_index()
    df_vw_previous = df_vw_filter[df_vw_filter['PR_Status'] == PR_Status_previous].reset_index()

    # 获得需要显示的年份数组，在图标中最为重要y轴坐标
    data_years = df_vw_filter.groupby(['YEAR']).size().reset_index().sort_values(['YEAR'], ascending=[True])['YEAR']
    # print(data_years)

    # 计算本轮和对比轮的总量
    df_vw_status = df_vw_filter.groupby(['PR_Status']).agg({'Volume': np.sum}).reset_index()
    df_vw_status_fuel_year = df_vw_filter.groupby(['PR_Status', 'Fuel_Type', 'YEAR']).agg(
        {'Volume': np.sum}).reset_index()
    PR_Status_local_vol = \
        df_vw_status[df_vw_status['PR_Status'] == PR_Status_local].reset_index().loc[0, 'Volume']
    PR_Status_previous_vol = \
        df_vw_status[df_vw_status['PR_Status'] == PR_Status_previous].reset_index().loc[0, 'Volume']
    print(PR_Status_local_vol)
    print(PR_Status_previous_vol)

    # 销售市场数据分组聚合，并按Fuel_type和year分组计算market effect rate(ice bev phev)-------------------------
    df_mkt_volume_local = \
        df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_local)].groupby(['Fuel_type', 'Year']) \
            .agg({'Volume': np.sum}).reset_index()
    print(df_mkt_volume_local)
    df_mkt_volume_previous = \
        df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_previous)].groupby(['Fuel_type', 'Year']) \
            .agg({'Volume': np.sum}).reset_index()
    print(df_mkt_volume_previous)

    year_mkt_rate_dict = OrderedDict()
    year_mkt_effect_dict = OrderedDict()
    year_fuel_mkt_effect_dict = OrderedDict()
    total_mkt_effect_volume = 0
    total_mkt_effect_percent = 0
    for year in data_years:
        fuel_mkt_rate_dict = {}
        fuel_mkt_effect_dict = {}
        fuel_mkt_effect_sum = 0
        for fuel in fuel_type_filter:
            mkt_volume_local = \
                df_mkt_volume_local[(df_mkt_volume_local['Year'] == year) & (df_mkt_volume_local['Fuel_type'] == fuel)] \
                    .reset_index().loc[0, 'Volume']
            mkt_volume_previous = \
                df_mkt_volume_previous[
                    (df_mkt_volume_local['Year'] == year) & (df_mkt_volume_local['Fuel_type'] == fuel)] \
                    .reset_index().loc[0, 'Volume']
            mkt_rate = mkt_volume_local / mkt_volume_previous - 1
            vw_volume_previous = df_vw_status_fuel_year[ \
                (df_vw_status_fuel_year['PR_Status'] == PR_Status_previous) & \
                (df_vw_status_fuel_year['YEAR'] == year) & \
                (df_vw_status_fuel_year['Fuel_Type'] == fuel)] \
                .reset_index().loc[0, 'Volume']
            mkt_effect_volume = vw_volume_previous * mkt_rate
            fuel_mkt_rate_dict[fuel] = mkt_rate
            fuel_mkt_effect_dict[fuel] = mkt_effect_volume
            fuel_mkt_effect_sum = fuel_mkt_effect_sum + mkt_effect_volume
        year_mkt_rate_dict[year] = fuel_mkt_rate_dict
        year_mkt_effect_dict[year] = fuel_mkt_effect_sum
        year_fuel_mkt_effect_dict[year] = fuel_mkt_effect_dict
        total_mkt_effect_volume = total_mkt_effect_volume + fuel_mkt_effect_sum
    total_mkt_effect_percent = total_mkt_effect_percent / PR_Status_previous_vol
    print('mkt_rate_dict=' + str(year_mkt_rate_dict))
    print('year_mkt_effect_dict=' + str(year_mkt_effect_dict))
    print('year_fuel_mkt_effect_dict' + str(year_fuel_mkt_effect_dict))
    print('total_mkt_effect_volume=' + str(total_mkt_effect_volume))

    # 计算Product effect
    # Product effect = New model family volume - Cancel model family volume
    vw_vol_local = df_vw_filter[df_vw_filter['PR_Status'] == PR_Status_local] \
        .groupby(['Model_Family', 'YEAR']).agg({'Volume': np.sum}).reset_index()
    vw_vol_previous = df_vw_filter[df_vw_filter['PR_Status'] == PR_Status_previous] \
        .groupby(['Model_Family', 'YEAR']).agg({'Volume': np.sum}).reset_index()
    model_familys = df_vw_filter.groupby(['Model_Family']).agg({'Volume': np.sum}) \
        .sort_values(['Volume'], ascending=[True]).reset_index()['Model_Family']
    print(model_familys)

    # 计算每year的每个model_family的New model family volume，如果不是new model则记为零
    # 计算每year的每个model_family的Cancel model family volume，如果不是Cancel model则记为零
    # 计算每year的Product effect = New model family volume - Cancel model family volume，如果没有则记为零
    new_model_vol_dict = OrderedDict()
    cancel_model_vol_dict = OrderedDict()
    product_effect_dict = OrderedDict()
    product_effect_total = 0
    for year in data_years:
        new_m_family_vol_dict = {}
        cancel_m_family_vol_dict = {}
        product_effect_year = 0
        new_m_family_vol_sum = 0
        cancel_m_family_vol_sum = 0
        for m_family in model_familys:
            m_family_local = 0
            m_family_previous = 0
            new_m_family_vol = 0
            cancel_m_family_vol = 0
            isEmpty_local = vw_vol_local[ \
                (vw_vol_local['YEAR'] == year) & (vw_vol_local['Model_Family'] == m_family)].empty
            isEmpty_previous = vw_vol_previous[ \
                (vw_vol_previous['YEAR'] == year) & (vw_vol_previous['Model_Family'] == m_family)].empty
            if not isEmpty_local:
                m_family_local = \
                    vw_vol_local[(vw_vol_local['YEAR'] == year) & (vw_vol_local['Model_Family'] == m_family)] \
                    .agg({'Volume': np.sum}).reset_index().iat[0, 1]
            if not isEmpty_previous:
                m_family_previous = \
                    vw_vol_previous[(vw_vol_previous['YEAR'] == year) & (vw_vol_previous['Model_Family'] == m_family)] \
                    .agg({'Volume': np.sum}).reset_index().iat[0, 1]
            if(m_family_local > 0 and m_family_previous <= 0):
                new_m_family_vol = m_family_local
            if(m_family_local <= 0 and m_family_previous > 0):
                cancel_m_family_vol = m_family_previous
            product_effect_m_family = new_m_family_vol - cancel_m_family_vol
            new_m_family_vol_dict[m_family] = new_m_family_vol
            cancel_m_family_vol_dict[m_family] = cancel_m_family_vol
            new_m_family_vol_sum = new_m_family_vol_sum + new_m_family_vol
            cancel_m_family_vol_sum = cancel_m_family_vol_sum + cancel_m_family_vol
        product_effect_year = new_m_family_vol_sum - cancel_m_family_vol_sum
        product_effect_total = product_effect_total + product_effect_year
        new_model_vol_dict[year] = new_m_family_vol_sum
        cancel_model_vol_dict[year] = cancel_m_family_vol_sum
        product_effect_dict[year] = product_effect_year
    product_effect_percent = product_effect_total / PR_Status_previous_vol
    print('product_effect_total = ' + str(product_effect_total))
    # print(new_model_vol_dict)
    # print(cancel_model_vol_dict)
    # print(product_effect_dict)

    # 开始计算Cycle Plan effect = SOP delay + SOP forward + EOP delay + EOP forward

    project_codes = df_vw_filter.groupby(['Project_Code']).agg({'Volume': np.sum}) \
        .sort_values(['Volume'], ascending=[True]).reset_index()['Project_Code']
    print(project_codes)


    def get_star_year(d):
        last_val = 0
        year = 0
        for key, val in d.items():
            if last_val == 0 and val != 0:
                year = key
                break
            else:
                last_val = val
        return year

    def get_end_year(d):
        last_val = 0
        year = 0
        for key, val in reversed(d.items()):
            if last_val == 0 and val != 0:
                year = key
                break
            else:
                last_val = val
        return year

    def get_PA_NF_status(code):
        is_num = code[-1].isdigit()
        status = 'NF'
        if is_num:
            status = 'PA'
        return status

    def sum_dict(a, b):
        c = OrderedDict()
        for key in a.keys() | b.keys():
            # c[key] = sum([d.get(key, 0) for d in (a, b)])
            c[key] = a.get(key, 0) + b.get(key, 0)
        return c


    # 1.计算SOP delay（负影响）
    total_SOP_delay_effect = 0
    year_SOP_delay_effect_dict = OrderedDict()
    year_mode_SOP_delay_effect_dict = OrderedDict()
    for code in project_codes:
        year_vol_local_dict = \
            df_vw_local[(df_vw_local['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        year_vol_previous_dict = \
            df_vw_previous[(df_vw_previous['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        star_year_local = get_star_year(year_vol_local_dict)
        star_year_previous = get_star_year(year_vol_previous_dict)
        # 判断当前的model是否为SOP delay，是否夸年
        mode_effect_dict = OrderedDict()
        mode_effect_sum = 0
        is_sop_delay = False
        is_cross_year = False
        if star_year_local > 0 and star_year_previous > 0 and star_year_local >= star_year_previous:
            if star_year_local == star_year_previous:
                sop_date_local = \
                    df_vw_local[(df_vw_local['Project_Code'] == code) & (df_vw_local['YEAR'] == star_year_local)] \
                        .reset_index().at[0, 'SOP_Time']
                sop_date_previous = \
                    df_vw_previous[
                        (df_vw_previous['Project_Code'] == code) & (df_vw_previous['YEAR'] == star_year_previous)] \
                        .reset_index().at[0, 'SOP_Time']
                if sop_date_local > sop_date_previous:
                    is_sop_delay = True
                    is_cross_year = False
                    # print('maoyadong-sop-delay#' + code + '#is_cross_years=false####' + \
                    # str(sop_date_local) + '#####' + str(sop_date_previous))
            else:
                is_sop_delay = True
                is_cross_year = True
                # print('maoyadong-sop-delay#' + code + '#is_cross_years=true####' + \
                # str(sop_date_local) + '#####' + str(sop_date_previous))

        if is_sop_delay:
            n = star_year_local - star_year_previous
            for year in data_years:
                mode_effect = 0
                if star_year_previous <= year <= star_year_local:
                    vw_previous_one_row = \
                        df_vw_previous[(df_vw_previous['Project_Code'] == code) & \
                                       (df_vw_previous['YEAR'] == year)].reset_index()
                    model_year_vol_previous = vw_previous_one_row.at[0, 'Volume']

                    if year == star_year_local:
                        fuel_type = vw_previous_one_row.at[0, 'Fuel_Type']
                        vw_local_one_row = \
                            df_vw_local[(df_vw_local['Project_Code'] == code) & \
                                        (df_vw_local['YEAR'] == year)].reset_index()
                        mkt_previous_vol = \
                            df_mkt_volume_previous[(df_mkt_volume_previous['Year'] == year) & \
                                                   (df_mkt_volume_previous['Fuel_type'] == fuel_type)] \
                                .reset_index().at[0, 'Volume']
                        mkt_rate_previous = model_year_vol_previous / mkt_previous_vol
                        previous_SOP_date = vw_previous_one_row.at[0, 'SOP_Time']
                        local_SOP_date = vw_local_one_row.at[0, 'SOP_Time']
                        if pd.isnull(previous_SOP_date) or pd.isnull(local_SOP_date):
                            print(previous_SOP_date)
                            print(local_SOP_date)
                        previous_SOP_month = int(previous_SOP_date.strftime("%m"))
                        local_SOP_month = int(local_SOP_date.strftime("%m"))
                        PA_NF_Status = get_PA_NF_status(code)

                        sale_previous_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'SOP') & \
                                                        (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                        (df_sale['Month'] == previous_SOP_month)].reset_index()
                        sale_local_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'SOP') & \
                                                     (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                     (df_sale['Month'] == local_SOP_month)].reset_index()

                        month_column_local = 'YEAR'
                        month_column_previous = 'YEAR'
                        rate_column = 'RATE_YEAR'
                        if is_cross_year:
                            if n > 0:
                                month_column_previous = month_column_previous + '+' + str(n)
                                month_column_local = month_column_local + '+' + str(n)
                                rate_column = rate_column + '+' + str(n)
                        else:
                            if previous_SOP_month >= 7:
                                month_column_previous = month_column_previous + '+1'
                            if local_SOP_month >= 7:
                                month_column_local = month_column_local + '+1'
                                rate_column = rate_column + '+1'

                        sale_month_previous = sale_previous_one_row.at[0, month_column_previous]
                        sale_month_local = sale_local_one_row.at[0, month_column_local]
                        sale_rate = sale_local_one_row.at[0, rate_column]  # .strip("%").astype(float)/100
                        mode_effect = \
                            mkt_rate_previous / sale_month_previous * sale_month_local * sale_rate \
                            * mkt_previous_vol - model_year_vol_previous

                        # if math.isnan(mode_effect):
                        #     print(code + '##' + str(year) + '##' + str(mkt_rate_previous) + '##' + \
                        #           str(sale_month_previous) + '##' + str(sale_month_local))
                    else:
                        mode_effect = 0 - model_year_vol_previous
                mode_effect_dict[year] = mode_effect
                mode_effect_sum = mode_effect_sum + mode_effect
            # print(mode_effect_dict)
            # print(mode_effect_sum)
            year_mode_SOP_delay_effect_dict[code] = mode_effect_dict
            total_SOP_delay_effect = total_SOP_delay_effect + mode_effect_sum
            # 将所有model相同年份的的EOP_delay_effect累计相加
            if len(year_SOP_delay_effect_dict) == 0:
                year_SOP_delay_effect_dict = mode_effect_dict
            elif len(mode_effect_dict) > 0:
                year_SOP_delay_effect_dict = sum_dict(year_SOP_delay_effect_dict, mode_effect_dict)
    print('year_SOP_delay_effect_dict='+str(year_SOP_delay_effect_dict))
    print('SOP delay total='+str(total_SOP_delay_effect))

    # 2.计算SOP forward（正影响）
    total_SOP_forward_effect = 0
    year_SOP_forward_effect_dict = OrderedDict()
    year_mode_SOP_forward_effect_dict = OrderedDict()
    for code in project_codes:
        year_vol_local_dict = \
            df_vw_local[(df_vw_local['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        year_vol_previous_dict = \
            df_vw_previous[(df_vw_previous['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        star_year_local = get_star_year(year_vol_local_dict)
        star_year_previous = get_star_year(year_vol_previous_dict)

        # 判断当前的model是否为SOP forward
        mode_effect_dict = OrderedDict()
        mode_effect_sum = 0
        is_sop_forward = False
        is_cross_year = False
        if star_year_local > 0 and star_year_previous > 0 and star_year_local <= star_year_previous:
            if star_year_local == star_year_previous:
                sop_date_local = \
                    df_vw_local[(df_vw_local['Project_Code'] == code) & (df_vw_local['YEAR'] == star_year_local)] \
                        .reset_index().at[0, 'SOP_Time']
                sop_date_previous = \
                    df_vw_previous[
                        (df_vw_previous['Project_Code'] == code) & (df_vw_previous['YEAR'] == star_year_previous)] \
                        .reset_index().at[0, 'SOP_Time']
                if sop_date_local < sop_date_previous:
                    is_sop_forward = True
                    is_cross_year = False
                    # print('maoyadong-sop-forward#' + code + '#is_cross_years=false####' + \
                    #       str(sop_date_local) + '#####' + str(sop_date_previous))
            else:
                is_sop_forward = True
                is_cross_year = True
                # print('maoyadong-sop-forward#' + code + '#is_cross_years=true####' + \
                #       str(sop_date_local) + '#####' + str(sop_date_previous))

        if is_sop_forward:
            n = star_year_previous - star_year_local
            i = 0
            for year in data_years:
                mode_effect = 0
                model_year_vol_previous = 0
                if star_year_local <= year <= star_year_previous:
                    rate_year = year + n
                    # print('maoyadong' + '#####' + str(n) + '###' + str(year) + '######' + \
                    #       str(rate_year) + '#####' + code + '####' + str(i))
                    if not df_vw_previous[(df_vw_previous['Project_Code'] == code) & (df_vw_previous['YEAR'] == year)] \
                            .empty:
                        model_year_vol_previous = \
                            df_vw_previous[(df_vw_previous['Project_Code'] == code) & (df_vw_previous['YEAR'] == year)] \
                                .reset_index().at[0, 'Volume']

                    vw_previous_rate_one_row = \
                        df_vw_previous[(df_vw_previous['Project_Code'] == code) & \
                                       (df_vw_previous['YEAR'] == rate_year)].reset_index()
                    model_year_rate_vol_previous = vw_previous_rate_one_row.at[0, 'Volume']

                    fuel_type = vw_previous_rate_one_row.at[0, 'Fuel_Type']

                    vw_local_one_row = \
                        df_vw_local[(df_vw_local['Project_Code'] == code) & \
                                    (df_vw_local['YEAR'] == year)].reset_index()

                    vw_local_rate_one_row = \
                        df_vw_local[(df_vw_local['Project_Code'] == code) & \
                                    (df_vw_local['YEAR'] == rate_year)].reset_index()
                    mkt_previous_vol = \
                        df_mkt_volume_previous[(df_mkt_volume_previous['Year'] == year) & \
                                               (df_mkt_volume_previous['Fuel_type'] == fuel_type)].reset_index().at[
                            0, 'Volume']
                    mkt_previous_rate_vol = \
                        df_mkt_volume_previous[(df_mkt_volume_previous['Year'] == rate_year) & \
                                               (df_mkt_volume_previous['Fuel_type'] == fuel_type)].reset_index().at[
                            0, 'Volume']

                    mkt_rate_previous = model_year_rate_vol_previous / mkt_previous_rate_vol

                    previous_SOP_date = vw_previous_rate_one_row.at[0, 'SOP_Time']
                    previous_SOP_month = int(previous_SOP_date.strftime("%m"))
                    local_SOP_date = vw_local_rate_one_row.at[0, 'SOP_Time']
                    local_SOP_month = int(local_SOP_date.strftime("%m"))
                    PA_NF_Status = get_PA_NF_status(code)
                    sale_previous_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'SOP') & \
                                                    (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                    (df_sale['Month'] == previous_SOP_month)].reset_index()
                    sale_local_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'SOP') & \
                                                 (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                 (df_sale['Month'] == local_SOP_month)].reset_index()

                    month_column_local = 'YEAR'
                    month_column_previous = 'YEAR'
                    rate_column = 'RATE_YEAR'
                    if is_cross_year:
                        if i > 0:
                            month_column_previous = month_column_previous + '+' + str(i)
                            month_column_local = month_column_local + '+' + str(i)
                            rate_column = rate_column + '+' + str(i)
                    else:
                        if previous_SOP_month >= 7:
                            month_column_previous = month_column_previous + '+1'
                        if local_SOP_month >= 7:
                            month_column_local = month_column_local + '+1'
                            rate_column = rate_column + '+1'

                    sale_month_previous = sale_previous_one_row.at[0, month_column_previous]
                    sale_month_local = sale_local_one_row.at[0, month_column_local]
                    sale_rate = sale_local_one_row.at[0, rate_column]  # .strip("%").astype(float)/100
                    mode_effect = \
                        mkt_rate_previous / sale_month_previous * sale_month_local * sale_rate \
                        * mkt_previous_vol - model_year_vol_previous
                    i = i + 1
                mode_effect_dict[year] = mode_effect
                mode_effect_sum = mode_effect_sum + mode_effect
            # print('mode_effect_dict'+str(mode_effect_dict))
            # print('mode_effect_sum='+str(mode_effect_sum))
            year_mode_SOP_forward_effect_dict[code] = mode_effect_dict
            total_SOP_forward_effect = total_SOP_forward_effect + mode_effect_sum
            # 将所有model相同年份的的EOP_delay_effect累计相加
            if len(year_SOP_forward_effect_dict) == 0:
                year_SOP_forward_effect_dict = mode_effect_dict
            elif len(mode_effect_dict) > 0:
                year_SOP_forward_effect_dict = sum_dict(year_SOP_forward_effect_dict, mode_effect_dict)
    print('year_SOP_forward_effect_dict=' + str(year_SOP_forward_effect_dict))
    print('SOP forward total='+str(total_SOP_forward_effect))

    # 3.计算EOP delay（正影响）
    total_EOP_delay_effect = 0
    year_EOP_delay_effect_dict = OrderedDict()
    year_mode_EOP_delay_effect_dict = OrderedDict()
    for code in project_codes:
        year_vol_local_dict = \
            df_vw_local[(df_vw_local['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        year_vol_previous_dict = \
            df_vw_previous[(df_vw_previous['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        end_year_local = get_end_year(year_vol_local_dict)
        end_year_previous = get_end_year(year_vol_previous_dict)

        mode_effect_dict = OrderedDict()
        mode_effect_sum = 0
        is_eop_delay = False
        is_cross_year = False
        # n = 0
        # 判断当前的model是否为EOP delay,且判断跨年或不跨年
        if end_year_local > 0 and end_year_previous > 0 and end_year_local >= end_year_previous:
            # print(str(star_year_local) + '#####' + str(star_year_previous))
            if end_year_local == end_year_previous:
                eop_date_local = \
                    df_vw_local[(df_vw_local['Project_Code'] == code) & (df_vw_local['YEAR'] == end_year_local)] \
                        .reset_index().at[0, 'EOP_Time']
                eop_date_previous = \
                    df_vw_previous[
                        (df_vw_previous['Project_Code'] == code) & (df_vw_previous['YEAR'] == end_year_previous)] \
                        .reset_index().at[0, 'EOP_Time']
                if eop_date_local > eop_date_previous:
                    is_eop_delay = True
                    is_cross_year = False
                    # n = 0
                    # print('maoyadong#' + code + '#EOP-delay is_cross_years=false####' + \
                    #       str(eop_date_local) + '#####' + str(eop_date_previous))
            else:
                is_eop_delay = True
                is_cross_year = True
                # print('maoyadong#' + code + '#EOP-delay is_cross_years=true####' + \
                #       str(end_year_local) + '#####' + str(end_year_previous))
                # n = end_year_local - end_year_previous

        if is_eop_delay:
            for year in data_years:
                # print(str(year) + "###" + model)
                mode_effect = 0
                if is_cross_year:
                    if end_year_previous < year <= end_year_local:
                        mode_effect = \
                            df_vw_local[(df_vw_local['Project_Code'] == code) & (df_vw_local['YEAR'] == year)] \
                            .reset_index().at[0, 'Volume']
                else:
                    if year == end_year_local:
                        vw_previous_one_row = \
                            df_vw_previous[(df_vw_previous['Project_Code'] == code) & \
                                           (df_vw_previous['YEAR'] == year)].reset_index()
                        model_year_vol_previous = vw_previous_one_row.at[0, 'Volume']

                        fuel_type = vw_previous_one_row.at[0, 'Fuel_Type']
                        vw_local_one_row = \
                            df_vw_local[(df_vw_local['Project_Code'] == code) & \
                                        (df_vw_local['YEAR'] == year)].reset_index()
                        mkt_previous_vol = \
                            df_mkt_volume_previous[(df_mkt_volume_previous['Year'] == year) & \
                                                   (df_mkt_volume_previous['Fuel_type'] == fuel_type)] \
                            .reset_index().at[0, 'Volume']

                        mkt_rate_previous = model_year_vol_previous / mkt_previous_vol
                        previous_EOP_date = vw_previous_one_row.at[0, 'EOP_Time']
                        previous_EOP_month = int(previous_EOP_date.strftime("%m"))
                        local_EOP_date = vw_local_one_row.at[0, 'EOP_Time']
                        local_EOP_month = int(local_EOP_date.strftime("%m"))
                        PA_NF_Status = get_PA_NF_status(code)
                        sale_previous_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'EOP') & \
                                                        (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                        (df_sale['Month'] == previous_EOP_month)].reset_index()
                        sale_local_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'EOP') & \
                                                     (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                     (df_sale['Month'] == local_EOP_month)].reset_index()
                        if previous_EOP_month <= 6:
                            month_column_previous = 'YEAR'
                        else:
                            month_column_previous = 'YEAR+1'
                        if local_EOP_month <= 6:
                            month_column_local = 'YEAR'
                        else:
                            month_column_local = 'YEAR+1'
                        sale_month_previous = sale_previous_one_row.at[0, month_column_previous]
                        sale_month_local = sale_local_one_row.at[0, month_column_local]

                        mode_effect = \
                            mkt_rate_previous / sale_month_previous * sale_month_local \
                            * mkt_previous_vol - model_year_vol_previous

                if math.isnan(mode_effect):
                    print(code+'##'+str(year)+'##'+str(mkt_rate_previous)+'##'+str(sale_month_previous)+ \
                          '##'+str(sale_month_local))
                mode_effect_dict[year] = mode_effect
                mode_effect_sum = mode_effect_sum + mode_effect
        # print('mode_effect_dict='+str(mode_effect_dict))
        # print('mode_effect_sum='+str(mode_effect_sum))

        year_mode_SOP_delay_effect_dict[code] = mode_effect_dict
        # 将所有model相同年份的的EOP_delay_effect累计相加
        if len(year_SOP_delay_effect_dict) == 0:
            year_EOP_delay_effect_dict = mode_effect_dict
        elif len(mode_effect_dict) > 0:
            year_EOP_delay_effect_dict = sum_dict(year_EOP_delay_effect_dict, mode_effect_dict)
        total_EOP_delay_effect = total_EOP_delay_effect + mode_effect_sum
    print('year_EOP_delay_effect_dict=' + str(year_EOP_delay_effect_dict))
    print('EOP delay total='+str(total_EOP_delay_effect))

    # 4.计算EOP forward（负影响）
    total_EOP_forward_effect = 0
    year_EOP_forward_effect_dict = OrderedDict()
    year_mode_EOP_forward_effect_dict = OrderedDict()
    for code in project_codes:
        year_vol_local_dict = \
            df_vw_local[(df_vw_local['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        year_vol_previous_dict = \
            df_vw_previous[(df_vw_previous['Project_Code'] == code)].sort_values(['YEAR'], ascending=[True]) \
                .reset_index().set_index('YEAR')['Volume'].to_dict()
        end_year_local = get_end_year(year_vol_local_dict)
        end_year_previous = get_end_year(year_vol_previous_dict)

        mode_effect_dict = OrderedDict()
        mode_effect_sum = 0
        is_eop_forward = False
        is_cross_year = False
        # n = 0
        # 判断当前的model是否为EOP delay,且判断跨年或不跨年
        if end_year_local > 0 and end_year_previous > 0 and end_year_local <= end_year_previous:
            # print(str(star_year_local) + '#####' + str(star_year_previous))
            if end_year_local == end_year_previous:
                eop_date_local = \
                    df_vw_local[(df_vw_local['Project_Code'] == code) & (df_vw_local['YEAR'] == end_year_local)] \
                        .reset_index().at[0, 'EOP_Time']
                eop_date_previous = \
                    df_vw_previous[
                        (df_vw_previous['Project_Code'] == code) & (df_vw_previous['YEAR'] == end_year_previous)] \
                        .reset_index().at[0, 'EOP_Time']
                if eop_date_local < eop_date_previous:
                    is_eop_forward = True
                    is_cross_year = False
                    # n = 0
                    # print('maoyadong#' + code + '#EOP-forward is_cross_years=false####' + \
                    #       str(eop_date_local) + '#####' + str(eop_date_previous))
            else:
                is_eop_forward = True
                is_cross_year = True
                # print('maoyadong#' + code + '#EOP-delay is_cross_years=true####' + \
                #       str(end_year_local) + '#####' + str(end_year_previous))
                # n = end_year_local - end_year_previous

        if is_eop_forward:
            for year in data_years:
                mode_effect = 0
                if is_cross_year:
                    if end_year_local < year <= end_year_previous:
                        mode_effect = \
                            df_vw_previous[(df_vw_previous['Project_Code'] == code) & (df_vw_previous['YEAR'] == year)] \
                            .reset_index().at[0, 'Volume']
                        mode_effect = 0 - mode_effect
                else:
                    if year == end_year_local:
                        vw_previous_one_row = \
                            df_vw_previous[(df_vw_previous['Project_Code'] == code) & \
                                           (df_vw_previous['YEAR'] == year)].reset_index()
                        model_year_vol_previous = vw_previous_one_row.at[0, 'Volume']

                        fuel_type = vw_previous_one_row.at[0, 'Fuel_Type']
                        vw_local_one_row = \
                            df_vw_local[(df_vw_local['Project_Code'] == code) & \
                                        (df_vw_local['YEAR'] == year)].reset_index()
                        mkt_previous_vol = \
                            df_mkt_volume_previous[(df_mkt_volume_previous['Year'] == year) & \
                                                   (df_mkt_volume_previous['Fuel_type'] == fuel_type)] \
                            .reset_index().at[0, 'Volume']

                        mkt_rate_previous = model_year_vol_previous / mkt_previous_vol
                        previous_EOP_date = vw_previous_one_row.at[0, 'EOP_Time']
                        previous_EOP_month = int(previous_EOP_date.strftime("%m"))
                        local_EOP_date = vw_local_one_row.at[0, 'EOP_Time']
                        local_EOP_month = int(local_EOP_date.strftime("%m"))

                        PA_NF_Status = get_PA_NF_status(code)

                        sale_previous_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'EOP') & \
                                                        (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                        (df_sale['Month'] == previous_EOP_month)].reset_index()
                        sale_local_one_row = df_sale[(df_sale['SOP_EOP_Status'] == 'EOP') & \
                                                     (df_sale['NF_PA_Status'] == PA_NF_Status) & \
                                                     (df_sale['Month'] == local_EOP_month)].reset_index()
                        if previous_EOP_month <= 6:
                            month_column_previous = 'YEAR'
                        else:
                            month_column_previous = 'YEAR+1'
                        if local_EOP_month <= 6:
                            month_column_local = 'YEAR'
                        else:
                            month_column_local = 'YEAR+1'
                        sale_month_previous = sale_previous_one_row.at[0, month_column_previous]
                        sale_month_local = sale_local_one_row.at[0, month_column_local]

                        mode_effect = \
                            mkt_rate_previous / sale_month_previous * sale_month_local \
                            * mkt_previous_vol - model_year_vol_previous

                if math.isnan(mode_effect):
                    print(code + '##' + str(year) + '##' + str(mkt_rate_previous) + '##' + \
                          str(sale_month_previous) + '##' + str(sale_month_local))
                mode_effect_dict[year] = mode_effect
                mode_effect_sum = mode_effect_sum + mode_effect
            # print('mode_effect_dict=' + str(mode_effect_dict))
            # print('mode_effect_sum=' + str(mode_effect_sum))

        year_mode_EOP_forward_effect_dict[code] = mode_effect_dict
        # 将所有model相同年份的的EOP_delay_effect累计相加
        if len(year_EOP_forward_effect_dict) == 0:
            year_EOP_forward_effect_dict = mode_effect_dict
        elif len(mode_effect_dict) > 0:
            year_EOP_forward_effect_dict = sum_dict(year_EOP_forward_effect_dict, mode_effect_dict)
        total_EOP_forward_effect = total_EOP_forward_effect + mode_effect_sum
    print('year_EOP_forward_effect_dict='+str(year_EOP_forward_effect_dict))
    print('EOP forward total='+str(total_EOP_forward_effect))

    # 合计Cycle Plan effect = SOP delay + SOP forward + EOP delay + EOP forward
    year_cycle_effect_dict = OrderedDict()
    total_cycle_plan_effect = 0
    for year in data_years:
        sop_delay_effect = year_SOP_delay_effect_dict[year]
        sop_forward_effect = year_SOP_forward_effect_dict[year]
        eop_delay_effect = year_EOP_delay_effect_dict[year]
        eop_forward_effect = year_EOP_forward_effect_dict[year]
        cycle_effect = sop_delay_effect + sop_forward_effect + eop_delay_effect + eop_forward_effect
        year_cycle_effect_dict[year] = cycle_effect
        total_cycle_plan_effect = total_cycle_plan_effect + cycle_effect

    total_cycle_plan_percent = total_cycle_plan_effect / PR_Status_previous_vol
    cycle_plan_effect = \
        total_SOP_delay_effect + total_SOP_forward_effect + total_EOP_delay_effect + total_EOP_forward_effect
    print('year_cycle_effect_dict=' + str(year_cycle_effect_dict))
    print('Total Cycle Plan effect=' + str(total_cycle_plan_effect))
    print('Cycle Plan effect='+str(cycle_plan_effect))

    # 5.计算Performance
    total_performance = 0
    year_performance_dict = OrderedDict()
    year_mode_performance_effect_dict = OrderedDict()
    for year in data_years:
        vol_local = \
            df_vw_local[(df_vw_local['YEAR'] == end_year_local)].reset_index().at[0, 'Volume']
        vol_previous = \
            df_vw_previous[(df_vw_previous['YEAR'] == end_year_local)].reset_index().at[0, 'Volume']
        mkt_effect = year_mkt_effect_dict[year]
        product_effect = product_effect_dict[year]
        cycle_effect = year_cycle_effect_dict[year]
        performance = vol_local - mkt_effect - product_effect - cycle_effect - vol_previous
        year_performance_dict[year] = performance
        total_performance = total_performance + performance
    total_performance_percent = total_performance / PR_Status_previous_vol
    print('year_performance_dict=' + str(year_performance_dict))
    print('total_performance=' + str(total_performance))

    effect_dict = OrderedDict()
    effect_percent_dict = OrderedDict()
    effect_year_dict = OrderedDict()

    effect_dict['Mkt Effect'] = total_mkt_effect_volume
    effect_percent_dict['Mkt Effect'] = total_mkt_effect_percent
    effect_year_dict['Mkt Effect'] = year_mkt_effect_dict

    effect_dict['Product Effect'] = product_effect_total
    effect_percent_dict['Product Effect'] = product_effect_percent
    effect_year_dict['Product Effect'] = product_effect_dict

    effect_dict['Cycle Plan Effect'] = total_cycle_plan_effect
    effect_percent_dict['Cycle Plan Effect'] = total_cycle_plan_percent
    effect_year_dict['Cycle Plan Effect'] = year_cycle_effect_dict

    effect_dict['Performance'] = total_performance
    effect_percent_dict['Performance'] = total_performance_percent
    effect_year_dict['Performance'] = year_performance_dict

    category_effect_list = ['Mkt Effect', 'Product Effect', 'Cycle Plan Effect', 'Performance']

    # 设置字体大小和起始位置
    top_base = 5
    font_size = Pt(10)

    # 设置堆积图y轴坐标和
    data_categories = []
    data_categories.insert(0, PR_Status_previous)
    data_categories.extend(category_effect_list)
    data_categories.append(PR_Status_local)
    print(data_categories)

    chart_data_stack = CategoryChartData()
    chart_data_stack.categories = data_categories

    # 设置柱状堆积图的series数据---------------------------------
    series_Volumes1 = []
    series_Volumes2 = []
    gap_sum = 0

    for key, effect in effect_dict.items():
        gap_sum += effect
        gap = PR_Status_previous_vol + gap_sum
        series_Volumes1.append(gap)
        series_Volumes2.append(abs(effect))

    series_Volumes1.insert(0, 0)
    series_Volumes1.append(0)

    series_Volumes2.insert(0, PR_Status_previous_vol)
    series_Volumes2.append(PR_Status_local_vol)

    series_Volumes2 = [vol / 1000 for vol in series_Volumes2]
    series_Volumes1 = [vol / 1000 for vol in series_Volumes1]
    chart_data_stack.add_series('series1', series_Volumes1)
    chart_data_stack.add_series('series2', series_Volumes2)
    print(series_Volumes1)
    print(series_Volumes2)

    x, y, cx, cy = Cm(1), Cm(top_base), Cm(24), Cm(8)
    graphic_frame_stack = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data_stack
    )

    chart_stack = graphic_frame_stack.chart
    # chart_stack.has_title = True
    # chart_stack.chart_title.has_text_frame = True
    # chart_stack.chart_title.text_frame.text = "maoyadong"
    # chart_stack.chart_title.text_frame.paragraphs[0].font.size = Pt(10)

    # chart_stack.has_legend = True
    # chart_stack.legend.position = XL_LEGEND_POSITION.LEFT  # XL_LEGEND_POSITION.CORNER
    # chart_stack.legend.include_in_layout = False
    # chart_stack.legend.font.size = Pt(10)

    chart_stack.series[0].data_labels.show_value = False
    chart_stack.series[0].format.fill.background()
    chart_stack.series[1].data_labels.show_value = True
    chart_stack.series[1].data_labels.number_format = '0'
    chart_stack.series[1].data_labels.font.size = font_size
    # chart_stack.series[1].data_labels.position = XL_LABEL_POSITION.ABOVE
    # for p in chart_stack.series[1].points:
    #     p.data_label.font.size = Pt(20)
    #     p.data_label.number_format = '0'

    for idx, category in enumerate(category_effect_list):
        index = idx + 1
        series_point = chart_stack.series[1].points[index]
        series_point.data_label.has_text_frame = True
        paragraphs1 = series_point.data_label.text_frame.paragraphs[0]
        paragraphs2 = series_point.data_label.text_frame.add_paragraph()
        run1 = paragraphs1.add_run()
        run1.text = format(effect_dict[category] / 1000, '.0f')
        run1.font.size = font_size
        run2 = paragraphs2.add_run()
        run2.text = format(effect_percent_dict[category], '.1%')
        run2.font.size = font_size


    # 设置中间柱子的颜色 负数为红 正数为绿
    for idx, series_point in enumerate(chart_stack.series[1].points):
        if 0 < idx < len(chart_stack.series[1].points) - 1:
            series_point.format.fill.patterned()  # 设置为可更改填充颜色的
            series_point.format.fill.solid()  # 设置为纯色填充
            if float(series_point.data_label.text_frame.paragraphs[0].runs[0].text) > 0:
                series_point.format.fill.fore_color.rgb = RGBColor(34, 139, 34)  # 绿色
            else:
                series_point.format.fill.fore_color.rgb = RGBColor(220, 20, 60)  # 红色

    value_axis_stack = chart_stack.value_axis
    value_axis_stack.has_major_gridlines = False
    value_axis_stack.major_tick_mark = XL_TICK_MARK.NONE
    value_axis_stack.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis_stack.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    value_axis_stack.visible = False

    category_axis_stack = chart_stack.category_axis
    category_axis_stack.has_major_gridlines = False
    category_axis_stack.major_tick_mark = XL_TICK_MARK.NONE
    category_axis_stack.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    category_axis_stack.tick_labels.font.size = font_size
    category_axis_stack.format.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
    category_axis_stack.visible = True

    # value_axis_stack.has_title = True
    # value_axis_stack.axis_title.has_text_frame = True
    # value_axis_stack.axis_title.text_frame.text = "False positive"
    # value_axis_stack.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

    # 开始创建表格
    rows = len(category_effect_list) + 1

    cols = len(data_years) + 1
    table_width = 24
    table_height = 3
    top = Cm(top_base + 8)
    left = Cm(1.5)  # Inches(2.0)
    width = Cm(table_width)  # Inches(6.0)
    height = Cm(table_height)  # Inches(0.8)

    # 添加表格到幻灯片 --------------------
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # 给data_year加E
    data_years_E = [str(year) + 'E' if year > start_year else str(year) for year in data_years]

    # 设置单元格宽度
    columns_width = table_width / cols - 0.1
    for i in range(cols):
        if i == 0:
            table.columns[i].width = Cm(columns_width + 0.6)
        else:
            table.columns[i].width = Cm(columns_width)  # Inches(2.0)

    row_height = table_height / rows
    for i in range(rows):
        table.rows[i].height = Cm(row_height)  # Inches(2.0)

    # 设置标题行
    for i in range(cols):
        if i == 0:
            table.cell(0, i).text = "vol.Change by Factors"
        else:
            table.cell(0, i).text = str(data_years_E[i - 1])


    # 填充表格第二行到最后一行数据
    for row_idx, category in enumerate(category_effect_list):
        table.cell(row_idx + 1, 0).text = category
        effect_dict = effect_year_dict[category]
        for col_idx, year in enumerate(data_years):
            effect_vol = effect_dict[year]
            table.cell(row_idx + 1, col_idx + 1).text = format(effect_vol / 1000, '.0f')


    # 调整table每个cell的字体
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(table):
        if cell.text.strip() == '':
            cell.text = r'/'
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = font_size

    # 添加ppt顶部的椭圆形shapes
    shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(12), Cm(top_base - 1), Cm(2), Cm(1.5))
    report_common.set_shape_oval_format(shape, font_size,PR_Status_local_vol - PR_Status_previous_vol, \
                                        (PR_Status_local_vol - PR_Status_previous_vol) / PR_Status_previous_vol)
    # shape.text_frame.word_wrap = False
    # shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    # paragraphs1 = shape.text_frame.paragraphs[0]
    # paragraphs2 = shape.text_frame.add_paragraph()
    # run1 = paragraphs1.add_run()
    # run1.text = format((PR_Status_local_vol - PR_Status_previous_vol) / 1000, '.0f')
    # run1.font.size = font_size
    # run1.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    # run2 = paragraphs2.add_run()
    # run2.text = format((PR_Status_local_vol - PR_Status_previous_vol) / PR_Status_previous_vol, '.1%')
    # run2.font.size = font_size
    # run2.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    # shape.fill.background()
    # shape.line.width = Pt(0.2)

    # 开始添加单位注释文本框
    # 在指定位置添加文本框
    textbox = shapes.add_textbox(Cm(1), Cm(top_base - 1), Cm(2), Cm(1))
    # 在文本框中写入文字
    textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = textbox.text_frame.add_paragraph()  # 新增段落
    run_t = para.add_run()
    run_t.text = "Volume '000units"  # 向段落写入文字
    run_t.font.size = Pt(8)


    # shape2 = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(13), Cm(1), Cm(1.5), Cm(1.5))
    # shape2.text = '2'
    # shape2.fill.background()
    # shape3 = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(26), Cm(4), Cm(1), Cm(1))
    # shape3.text = '3'
    # shape3.fill.background()
    #
    # connector = shapes.add_connector(
    #     MSO_CONNECTOR.ELBOW, 1, 1, 1, 1
    # )
    # connector2 = shapes.add_connector(
    #     MSO_CONNECTOR.ELBOW, 1, 1, 1, 1
    # )
    # connector.begin_connect(shape, 0)
    # connector.end_connect(shape2, 1)
    # connector2.begin_connect(shape2, 3)
    # connector2.end_connect(shape3, 0)

    prs.save('c:/auto-report/template_tmp6.pptx')
    print("maoyadong")
