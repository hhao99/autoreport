from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx import Presentation
import pandas as pd
from os.path import dirname, join

class ReportBase(object):
    def __init__(self, reportId):
        self.id = reportId

    def LoadConfig(self, config):
        # open ppt with cover
        # self.df_vw = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 2)
        # self.df_mkt = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 3)
        database_small_demo = join(dirname(__file__), '../sample/Database_small_demo.xlsx')

        self.vw_path = r'c:/auto-report/Database_small_demo.xlsx'
        self.mkt_path = r'c:/auto-report/Database_small_demo.xlsx'
        self.template_path = r'c:/auto-report/cover.pptx'
        self.save_path= r'c:/auto-report/template_tmp'
        self.df_vw = pd.read_excel(self.vw_path, 2)
        self.df_mkt = pd.read_excel(self.mkt_path, 3)
        oem = self.df_vw.groupby('OEM').size().reset_index()["OEM"]
        brand = self.df_vw.groupby('Brand').size().reset_index()["Brand"]
        build_type = self.df_vw.groupby('Build_Type').size().reset_index()["Build_Type"]
        fuel_type = self.df_vw.groupby('Fuel_Type').size().reset_index()["Fuel_Type"]
        fuel_type_group = self.df_vw.groupby('Fuel_Type_Group').size().reset_index()["Fuel_Type_Group"]
        print(oem)
        print(brand)
        print(build_type)
        print(fuel_type)
        print(fuel_type_group)

        # self.PR_Status_local = config.pr_state[0]
        # self.PR_Status_previous = config.pr_state_2[0]
        # self.year_expect_filter = config.YEAR
        # self.year_actual_filter = config.YEAR_2
        self.PR_Status_actual = 'Actual'
        self.PR_Status_local = 'PR66.OP'
        self.PR_Status_previous = 'PR66.SP'
        self.year_expect_filter = [2018, 2019, 2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027]
        self.year_actual_filter = [2016, 2017]

        self.PR_Status = []
        self.year_filter = []
        self.year_categories = []
        if self.id in [0, 1, 2, 3]:
            self.PR_Status.append(self.PR_Status_actual)
            self.PR_Status.append(self.PR_Status_local)
            self.year_filter = self.year_actual_filter + self.year_expect_filter
            self.year_categories = get_year_categories(self.year_actual_filter, self.year_expect_filter)
        else:
            self.PR_Status.append(self.PR_Status_local)
            self.PR_Status.append(self.PR_Status_previous)
            self.year_filter = self.year_expect_filter
            self.year_categories = get_year_categories([], self.year_expect_filter)

        self.year_filter.sort()

        # self.oem_filter = config.OEM
        # self.brand_filter = config.Brand
        # self.build_type_filter = config.Build_Type
        # self.all_mkt = config.All_mkt
        # self.fuel_type_filter = config.Fuel_Type
        # self.view_rules = config.group_by
        self.oem_filter = oem
        self.brand_filter = brand
        self.build_type_filter = build_type
        self.fuel_type_all = fuel_type
        self.fuel_type_filter = fuel_type
        self.fuel_type_group_all = fuel_type_group
        self.fuel_type_group_filter = fuel_type_group
        self.view_rules = 'Brand'
        self.all_mkt = False

    def LoadBaseData(self):
        self.df_vw = pd.read_excel(self.vw_path, 2)
        self.df_mkt = pd.read_excel(self.mkt_path, 3)

    def FilterData(self):
        self.df_vw_filter = \
            self.df_vw[(self.df_vw['PR_Status'].isin(self.PR_Status)) & (self.df_vw['YEAR'].isin(self.year_filter)) & ( \
            self.df_vw['OEM'].isin(self.oem_filter)) & (self.df_vw['Brand'].isin(self.brand_filter)) & ( \
            self.df_vw['Build_Type'].isin(self.build_type_filter))]
        self.df_mkt_filter = \
            self.df_mkt[(self.df_mkt['Status'].isin(self.PR_Status)) & (self.df_mkt['Year'].isin(self.year_filter))]
        print("df_vw_filter=" + str(self.df_vw_filter.shape[0]))
        print("df_mkt_filter=" + str(self.df_mkt_filter.shape[0]))

    def CreateSlide(self):
        self.shapes, self.prs = get_shapes(self.template_path)

    def SaveSlide(self, reportid=0):
        full_path = self.save_path+str(reportid)+'.pptx'
        self.prs.save(full_path)
        print("Save ppt success = " + full_path)

    def Run(self):
        print('Call report base')

class Report(object):
    def __init__(self, name, score):
        self.name = name
        self.score = score

class PrivateConfig(object):
    def __init__(self, cofig):
        pass



def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


# 设置table 字体和格式
def set_table_format(table, font_size, content):
    for cell in iter_cells(table):
        if cell.text.strip() == '':
            cell.text = content
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = font_size


# 设置textbox的字体和格式
def set_textbox_format(textbox, font_size, content):
    textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = textbox.text_frame.add_paragraph()  # 新增段落
    run_t = para.add_run()
    run_t.text = content  # 向段落写入文字
    run_t.font.size = font_size


def set_shape_oval_format(shape, font_size, volume, rate):
    shape.text_frame.word_wrap = False
    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraphs1 = shape.text_frame.paragraphs[0]
    paragraphs2 = shape.text_frame.add_paragraph()
    run1 = paragraphs1.add_run()
    run1.text = format(volume, '.0f')
    run1.font.size = font_size
    run1.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    run2 = paragraphs2.add_run()
    run2.text = format(rate, '.1%')
    run2.font.size = font_size
    run2.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    shape.fill.background()
    shape.line.width = Pt(0.2)


def get_shapes(template_path, templateid=1):
    prs = Presentation(template_path)
    title_only_slide_layout = prs.slide_layouts[templateid]
    slide = prs.slides.add_slide(title_only_slide_layout)
    return slide.shapes, prs


def get_fuel_from_fuelTypeGroup(fuelTypeGroup):
    if fuelTypeGroup == 'ICE':
        fuel = ['ICE']
    elif fuelTypeGroup == 'NEV':
        fuel = ['BEV', 'PHEV']
    else:
        fuel = ['ICE']
    return fuel


def get_fuel_list():
    return ['ICE', 'BEV', 'PHEV']


def get_fuel_group_list():
    return ['ICE', 'NEV']


def get_year_categories(year_actual_list, year_expect_lsit):
    year_actual_list.sort()
    year_expect_lsit.sort()
    return [str(year) for year in year_actual_list] + [str(year) + 'E' for year in year_expect_lsit]
