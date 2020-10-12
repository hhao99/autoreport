from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE, XL_LABEL_POSITION
from pptx.util import Cm  # Inches
from pandas import np
from pptx.enum.chart import XL_LEGEND_POSITION

from pptx.dml.color import RGBColor
import pandas as pd

if __name__ == '__main__':

    df_vw = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 0)
    df_mkt = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 1)
    prs = Presentation('c:/auto-report/cover.pptx')
    title_only_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    print("df_vw=" + str(df_vw.shape[0]))
    print("df_mkt=" + str(df_mkt.shape[0]))
    view_rules = 'OEM'  # or OEM
    all_mkt = False  # true分母是否为全市场，false分母为fuel_type细分市场
    oem = df_vw.groupby('OEM').size().reset_index()["OEM"]
    brand = df_vw.groupby('Brand').size().reset_index()["Brand"]
    fuel_type = df_vw.groupby('Fuel_Type').size().reset_index()["Fuel_Type"]
    fuel_type_group = df_vw.groupby('Fuel_Type_Group').size().reset_index()["Fuel_Type_Group"]
    Platform_Group = df_vw.groupby('Platform_Group').size().reset_index()["Platform_Group"]
    print(oem)
    print(brand)
    print(fuel_type)
    print(fuel_type_group)
    print(Platform_Group)

    # 此处应该对 oem brand fuel_type三个数组进行过滤，把用户不需要的删除掉，默认是全选，需要通过读取web端的配置文件
    oem_filter = oem
    brand_filter = brand
    Platform_Group_filter = ['BEV', 'FBU', 'MEB', 'MQB']#Platform_Group
    fuel_type_filter = fuel_type  # 根据私有fitler决定哪个fuel_type需要保留
    fuel_type_group_filter = 'NEV' #在这报表中，该过滤选项是单选

    #0BEV1FBU2MEB3MLB4MQB5PPE6PQ

    if fuel_type_group_filter == 'ICE':
        mkt_fuel_filter = ['ICE']
    else:
        mkt_fuel_filter = ['PHEV', 'BEV']

    view_rules = 'Brand'
    PR_Status_local = 'PR67.OP'  # 本轮
    PR_Status = []
    PR_Status.append(PR_Status_local)
    start_year = 2018
    year_span = 9
    end_year = start_year + year_span

    # 根据所有filter过滤大众数据
    df_vw_filter = \
        df_vw[(df_vw['PR_Status'].isin(PR_Status)) & (df_vw['YEAR'] >= start_year) & (df_vw['YEAR'] <= end_year) \
              & (df_vw['OEM'].isin(oem_filter)) & (df_vw['Brand'].isin(brand_filter)) \
              & (df_vw['Fuel_Type_Group'] == fuel_type_group_filter) \
              & (df_vw['Platform_Group'].isin(Platform_Group_filter))]
    print("df_vw_filter=" + str(df_vw_filter.shape[0]))

    # 根据所有filter过滤大众数据,不在所选平台内的数据
    df_vw_filter_notin_Platform = \
        df_vw[(df_vw['PR_Status'].isin(PR_Status)) & (df_vw['YEAR'] >= start_year) & (df_vw['YEAR'] <= end_year) \
              & (df_vw['OEM'].isin(oem_filter)) & (df_vw['Brand'].isin(brand_filter)) \
              & (df_vw['Fuel_Type_Group'] == fuel_type_group_filter) \
              & (-df_vw['Platform_Group'].isin(Platform_Group_filter))]
    print("df_vw_filter_notin_Platform=" + str(df_vw_filter_notin_Platform.shape[0]))

    df_mkt_filter = \
        df_mkt[(df_mkt['Status'].isin(PR_Status)) & (df_mkt['Year'] >= start_year) & (df_mkt['Year'] <= end_year) \
              & (df_mkt['Fuel_type'].isin(fuel_type))]
    print("df_mkt_filter=" + str(df_mkt_filter.shape[0]))

    # 获得需要显示的年份数组，在图标中最为重要y轴坐标
    data_years = df_vw_filter.groupby(['YEAR']).size().reset_index().sort_values(['YEAR'], ascending=[True])['YEAR']
    # print(data_years)

    #按照降序排序出所有需要显示的Platform_Group
    platform_group_order = df_vw_filter.groupby(['Platform_Group']).agg({'Volume': np.sum}).sort_values(['Volume'], ascending=[False]).reset_index()['Platform_Group']
    # print(fuel_type_order)

    #根据全市场还是细分市场 获得mkt表的每一年的销量
    mkt_year_dict = {}
    if all_mkt:
        mkt_year_dict = \
            df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_local) & (df_mkt_filter['Fuel_type'].isin(fuel_type))] \
            .groupby('Year').agg({'Volume': np.sum}).sort_values(['Year'], ascending=[True]).reset_index().set_index('Year')['Volume'].to_dict()
    else:
        mkt_year_dict = \
            df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_local) & (df_mkt_filter['Fuel_type'].isin(mkt_fuel_filter))] \
            .groupby('Year').agg({'Volume': np.sum}).sort_values(['Year'], ascending=[True]).reset_index().set_index('Year')['Volume'].to_dict()
    print(mkt_year_dict)

    # 计算本轮每年的选中的fuel_type_group的量(只统计需要显示的)
    year_fuel_group_volumes = df_vw_filter.groupby(['Fuel_Type_Group', 'YEAR']).agg({'Volume': np.sum}).sort_values(['YEAR'], ascending=[True]).reset_index()['Volume']
    print(year_fuel_group_volumes)

    # 计算每年每个fuel_type_group的MS%
    ms_year_fuel_group_list = []
    for idx, year in enumerate(data_years):
        vw_volume = year_fuel_group_volumes[idx]
        mkt_volume = mkt_year_dict[year]
        ms = vw_volume / mkt_volume
        ms_year_fuel_group_list.append(ms)
    print(ms_year_fuel_group_list)

    #计算每个Platform_Group本轮每年的volume
    df_vw_platform_year_vol = df_vw_filter.groupby(['YEAR', 'Platform_Group']).agg({'Volume': np.sum}).reset_index()
    year_platform_vol = {}
    for plat in platform_group_order:
        vol_list = []
        for idx, year in enumerate(data_years):
            vol_local = 0
            if not df_vw_platform_year_vol[(df_vw_platform_year_vol['Platform_Group'] == plat) & (df_vw_platform_year_vol['YEAR'] == year)].empty:
                vol_local = \
                    df_vw_platform_year_vol[(df_vw_platform_year_vol['Platform_Group'] == plat) & (df_vw_platform_year_vol['YEAR'] == year)].reset_index().loc[0, 'Volume']
            vol_list.append(vol_local)
        year_platform_vol[plat] = vol_list
    print(year_platform_vol)


    # 计算在选择的Platform_Group之外的车，按ICE或NEV进行数量统计
    year_not_platform_vol = {}
    if fuel_type_group_filter == 'ICE':
        year_other_vol = df_vw_filter_notin_Platform.groupby(['YEAR']).agg({'Volume': np.sum}).sort_values(['YEAR'], ascending=[True]).reset_index()['Volume']
        year_not_platform_vol['OTHER'] = year_other_vol
    else:
        year_phev_vol = df_vw_filter_notin_Platform[df_vw_filter_notin_Platform['Fuel_Type'] == 'PHEV'].groupby(['YEAR']).agg({'Volume': np.sum}).sort_values(['YEAR'], ascending=[True]).reset_index()['Volume']
        year_other_nev_vol = df_vw_filter_notin_Platform[df_vw_filter_notin_Platform['Fuel_Type'] == 'BEV'].groupby(['YEAR']).agg({'Volume': np.sum}).sort_values(['YEAR'], ascending=[True]).reset_index()['Volume']
        year_not_platform_vol['PHEV'] = year_phev_vol
        year_not_platform_vol['OTHER NEV'] = year_other_nev_vol

    #开始绘制柱状图------------------------------------
    # 设置柱状图图y轴坐标和
    chart_data_cluster = CategoryChartData()
    chart_data_cluster.categories = data_years

    # 设置柱状图的series数据--------------------------------
    series_Volumes1 = [vol / 1000 for vol in year_fuel_group_volumes]
    print(year_fuel_group_volumes)
    chart_data_cluster.add_series(fuel_type_group_filter + ' VOL', series_Volumes1)
    #chart_data_cluster.add_series(fuel_type_group_filter + ' VOL2', series_Volumes1)

    x, y, cx, cy = Cm(1), Cm(7.3), Cm(24), Cm(5)
    graphic_frame_cluster = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_cluster
    )

    chart_cluster = graphic_frame_cluster.chart
    # chart_stack.has_title = True
    # chart_stack.chart_title.has_text_frame = True
    # chart_stack.chart_title.text_frame.text = "maoyadong"
    # chart_stack.chart_title.text_frame.paragraphs[0].font.size = Pt(10)

    chart_cluster.has_legend = True
    chart_cluster.legend.position = XL_LEGEND_POSITION.LEFT  # XL_LEGEND_POSITION.CORNER
    chart_cluster.legend.include_in_layout = False
    chart_cluster.legend.font.size = Pt(10)

    for cluster_serie in chart_cluster.series:
        cluster_serie.data_labels.show_value = True
        cluster_serie.data_labels.number_format = '0'
        cluster_serie.data_labels.font.size = Pt(8)
        # stack_serie.data_labels.position = XL_DATA_LABEL_POSITION.ABOVE

    value_axis_cluster = chart_cluster.value_axis
    value_axis_cluster.has_major_gridlines = False
    value_axis_cluster.major_tick_mark = XL_TICK_MARK.NONE
    value_axis_cluster.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis_cluster.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    value_axis_cluster.visible = False

    category_axis_cluster = chart_cluster.category_axis
    category_axis_cluster.has_major_gridlines = False
    category_axis_cluster.major_tick_mark = XL_TICK_MARK.NONE
    category_axis_cluster.tick_label_position = XL_TICK_LABEL_POSITION.LOW #NONE
    category_axis_cluster.tick_labels.font.size = Pt(8)
    category_axis_cluster.format.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
    category_axis_cluster.visible = True

    # 开始创建点线图-----------------------------------
    chart_data_line = ChartData()
    chart_data_line.categories = data_years

    # 设置折线图的series数据--------------------------------
    print(ms_year_fuel_group_list)
    chart_data_line.add_series(fuel_type_group_filter + ' MS%', ms_year_fuel_group_list)
    #chart_data_line.add_series(fuel_type_group_filter + ' MS%2', ms_year_fuel_group_list)

    x, y, cx, cy = Cm(1), Cm(5.5), Cm(24), Cm(3)
    chart_line = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_line
    ).chart

    chart_line.has_legend = True
    chart_line.legend.include_in_layout = False
    chart_line.legend.position = XL_LEGEND_POSITION.LEFT
    chart_line.legend.font.size = Pt(7)

    for line_serie in chart_line.series:
        line_serie.smooth = True
        line_serie.marker.style = XL_MARKER_STYLE.CIRCLE
        line_serie.data_labels.show_value = True
        line_serie.data_labels.number_format = '0.00%'
        line_serie.data_labels.font.size = Pt(8)

    chart_line.series[0].data_labels.position = XL_LABEL_POSITION.ABOVE
    #chart_line.series[1].data_labels.position = XL_LABEL_POSITION.BELOW

    value_axis_line = chart_line.value_axis
    value_axis_line.has_major_gridlines = False
    value_axis_line.major_tick_mark = XL_TICK_MARK.NONE
    value_axis_line.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis_line.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    value_axis_line.visible = False

    category_axis_line = chart_line.category_axis
    category_axis_line.has_major_gridlines = False
    category_axis_line.major_tick_mark = XL_TICK_MARK.NONE
    category_axis_line.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    category_axis_line.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    category_axis_line.visible = False


    # 开始创建表格
    #计算表格的行数和列数
    rows = len(platform_group_order) + len(year_not_platform_vol) + 1
    cols = len(data_years) + 1

    table_width = 24
    table_height = 2
    top = Cm(12.3)
    left = Cm(1.5)  # Inches(2.0)
    width = Cm(table_width)  # Inches(6.0)
    height = Cm(table_height)  # Inches(0.8)

    # 添加表格到幻灯片 --------------------
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # 给data_year加E
    data_years_E = [str(year) + 'E' if year > start_year else str(year) for year in data_years]

    # 设置单元格宽度
    columns_width = table_width / cols - 0.1
    for i in range(cols):
        if i == 0:
            table.columns[i].width = Cm(columns_width + 0.3)
        else:
            table.columns[i].width = Cm(columns_width)  # Inches(2.0)

    row_height = table_height / rows
    for i in range(rows):
        table.rows[i].height = Cm(row_height)  # Inches(2.0)

    # 设置标题行
    for i in range(cols):
        if i == 0:
            table.cell(0, i).text = "Vol.By Type " + view_rules
        else:
            table.cell(0, i).text = str(data_years_E[i - 1])

    # 设置第一行和第二行数据内容
    row_idx = 1
    for key, val in year_not_platform_vol.items():
        table.cell(row_idx, 0).text = key
        for col_idx, vol in enumerate(val):
            if col_idx > 0:
                table.cell(row_idx, col_idx).text = format(vol / 1000, '.0f')
        row_idx = row_idx + 1

    # 填充表格第三行（或第二行）到最后一行数据
    start_row_index = len(year_not_platform_vol) + 1
    for idx, plat in enumerate(platform_group_order):
        row_idx = idx + start_row_index
        table.cell(row_idx, 0).text = plat
        vol_list = year_platform_vol[plat]
        for col_idx, vol in enumerate(vol_list):
            table.cell(row_idx, col_idx + 1).text = format(vol, '.0f')


    # 调整table每个cell的字体
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(table):
        if cell.text.strip() == '':
            cell.text = r'/'
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(8)

    # 开始添加注释文本框
    left = Cm(1)  # left，top为相对位置
    top = Cm(4)
    width = Cm(2)  # width，height为文本框的大小
    height = Cm(1)

    # 在指定位置添加文本框
    textbox = shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    # 在文本框中写入文字
    para = tf.add_paragraph()  # 新增段落
    para.text = "Volume '000units"  # 向段落写入文字
    para.line_spacing = 1.5  # 1.5 倍的行距
    para.font.size = Pt(6)

    prs.save('c:/auto-report/template_tmp4.pptx')
    print("maoyadong")
