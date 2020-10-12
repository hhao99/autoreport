from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE, XL_DATA_LABEL_POSITION
from pptx.util import Cm  # Inches
from pandas import np
from pptx.enum.chart import XL_LEGEND_POSITION

from pptx.dml.color import RGBColor
import pandas as pd

if __name__ == '__main__':

    df_vw = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 0)
    df_mkt = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', 1)
    prs = Presentation('c:/auto-report/cover.pptx')
    title_only_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    print("df_vw=" + str(df_vw.shape[0]))
    print("df_mkt=" + str(df_mkt.shape[0]))
    view_rules = 'OEM'  # or OEM
    all_mkt = False  # true分母是否为全市场，false分母为fuel_type细分市场
    oem = df_vw.groupby('OEM').size().reset_index()["OEM"]
    brand = df_vw.groupby('Brand').size().reset_index()["Brand"]
    fuel_type = df_vw.groupby('Fuel_Type').size().reset_index()["Fuel_Type"]
    fuel_type_group = df_vw.groupby('Fuel_Type_Group').size().reset_index()["Fuel_Type_Group"]

    print(oem)
    print(brand)
    print(fuel_type)
    print(fuel_type_group)
    # 此处应该对 oem brand fuel_type三个数组进行过滤，把用户不需要的删除掉，默认是全选，需要通过读取web端的配置文件
    oem_filter = oem
    brand_filter = brand
    fuel_type_filter = fuel_type  # 根据私有fitler决定哪个fuel_type需要保留

    view_rules = 'Brand'
    PR_Status_local = 'PR67.OP'  # 本轮
    PR_Status_previous = 'PR67.SP'  # 对比轮也叫上一轮
    PR_Status = []
    PR_Status.append(PR_Status_local)
    PR_Status.append(PR_Status_previous)
    start_year = 2018
    year_span = 9
    end_year = start_year + year_span

    # 根据所有filter过滤大众数据
    df_vw_filter = \
        df_vw[(df_vw['PR_Status'].isin(PR_Status)) & (df_vw['YEAR'] >= start_year) & (df_vw['YEAR'] <= end_year) \
              & (df_vw['OEM'].isin(oem)) & (df_vw['Brand'].isin(brand)) & (df_vw['Fuel_Type'].isin(fuel_type))]
    print("df_vw_filter=" + str(df_vw_filter.shape[0]))

    df_mkt_filter = \
        df_mkt[(df_mkt['Status'].isin(PR_Status)) & (df_mkt['Year'] >= start_year) & (df_mkt['Year'] <= end_year) \
              & (df_mkt['Fuel_type'].isin(fuel_type))]
    print("df_mkt_filter=" + str(df_mkt_filter.shape[0]))

    # 获得需要显示的年份数组，在图标中最为重要y轴坐标
    data_years = df_vw_filter.groupby(['YEAR']).size().reset_index().sort_values(['YEAR'], ascending=[True])['YEAR']
    # print(data_years)

    #根据全市场还是细分市场 获得mkt表的每一年的销量
    df_mkt_year_local = []
    df_mkt_year_previous = []
    if all_mkt:
        df_mkt_year_local = \
            df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_local) & (df_mkt_filter['Fuel_type'].isin(fuel_type))] \
            .groupby('Year').agg({'Volume': np.sum}).sort_values(['Year'], ascending=[True]).reset_index()['Volume']
        df_mkt_year_previous = \
            df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_previous) & (df_mkt_filter['Fuel_type'].isin(fuel_type))] \
            .groupby('Year').agg({'Volume': np.sum}).sort_values(['Year'], ascending=[True]).reset_index()['Volume']
    else:
        df_mkt_year_local = df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_local) & (df_mkt_filter['Fuel_type'].isin(fuel_type_filter))] \
            .groupby('Year').agg({'Volume': np.sum}).sort_values(['Year'], ascending=[True]).reset_index()['Volume']
        df_mkt_year_previous = \
            df_mkt_filter[(df_mkt_filter['Status'] == PR_Status_previous) & (df_mkt_filter['Fuel_type'].isin(fuel_type_filter))] \
            .groupby('Year').agg({'Volume': np.sum}).sort_values(['Year'], ascending=[True]).reset_index()['Volume']
    print(df_mkt_year_local)
    print(df_mkt_year_previous)

    # 计算本轮和对比轮的每年的量
    df_vw_group_status_year = df_vw_filter.groupby(['PR_Status', 'YEAR']).agg({'Volume': np.sum}).reset_index()
    year_volumes_local = df_vw_group_status_year[(df_vw_group_status_year['PR_Status'] == PR_Status_local)].sort_values(['YEAR'], ascending=[True]).reset_index()['Volume']
    year_volumes_previous = df_vw_group_status_year[(df_vw_group_status_year['PR_Status'] == PR_Status_previous)].sort_values(['YEAR'], ascending=[True]).reset_index()['Volume']
    print(year_volumes_local)
    print(year_volumes_previous)

    # 计算每年的 本轮和对比轮的总的MS% 和 MS% change
    year_ms_local = []
    year_ms_previous = []
    year_ms_change = []
    for index in range(len(data_years)):
        ms_local = year_volumes_local[index] / df_mkt_year_local[index]
        ms_previous = year_volumes_previous[index] / df_mkt_year_previous[index]
        ms_change = ms_local - ms_previous
        year_ms_local.append(ms_local)
        year_ms_previous.append(ms_previous)
        year_ms_change.append(ms_change)
    print(year_ms_local)
    print(year_ms_previous)
    print(year_ms_change)

    # 按大众本轮的销量排序OEM和Brand
    categories_order = \
        df_vw_filter[df_vw_filter['PR_Status'].isin(PR_Status)].groupby(view_rules).agg({'Volume': np.sum}) \
        .reset_index().sort_values(['Volume'], ascending=[False]).reset_index()[view_rules]
    print(categories_order)

    # 计算本轮和对比轮的每个Brand（或OEM）每年的量
    print(df_vw_filter.groupby([view_rules]).agg({'Volume': np.sum}).reset_index())
    df_vw_group_status_year_cat = df_vw_filter.groupby(['PR_Status', 'YEAR', view_rules]).agg({'Volume': np.sum}).reset_index()
    df_year_category_local = df_vw_group_status_year_cat[(df_vw_group_status_year_cat['PR_Status'] == PR_Status_local)].reset_index()
    df_year_category_previous = df_vw_group_status_year_cat[(df_vw_group_status_year_cat['PR_Status'] == PR_Status_previous)].reset_index()
    print(df_year_category_local)
    print(df_year_category_previous)

    #计算每个brand（或OEM）每年的 本轮和对比轮的总的MS% 和 MS% change
    year_category_ms_change = {}
    for category in categories_order:
        ms_change_list = []
        for idx, year in enumerate(data_years):
            vol_local = 0
            vol_previous = 0
            if not df_year_category_local[(df_year_category_local[view_rules] == category) & (df_year_category_local['YEAR'] == year)].empty:
                vol_local = \
                    df_year_category_local[(df_year_category_local[view_rules] == category) & (df_year_category_local['YEAR'] == year)].reset_index().loc[0, 'Volume']
            if not df_year_category_previous[(df_year_category_previous[view_rules] == category) & (df_year_category_previous['YEAR'] == year)].empty:
                vol_previous = \
                    df_year_category_previous[(df_year_category_previous[view_rules] == category) & (df_year_category_previous['YEAR'] == year)].reset_index().loc[0, 'Volume']
            ms_change = vol_local / df_mkt_year_local[idx] - vol_previous / df_mkt_year_previous[idx]
            ms_change_list.append(ms_change)
        year_category_ms_change[category] = ms_change_list
    print(year_category_ms_change)

    # 开始创建点线图-----------------------------------
    chart_data_line = ChartData()
    chart_data_line.categories = data_years

    series_line_Volumes1 = year_ms_previous
    series_line_Volumes2 = year_ms_local
    chart_data_line.add_series(PR_Status_previous, series_line_Volumes1)
    chart_data_line.add_series(PR_Status_local, series_line_Volumes2)

    x, y, cx, cy = Cm(1), Cm(4.5), Cm(24), Cm(4)
    chart_line = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_line
    ).chart

    chart_line.has_legend = True
    chart_line.legend.include_in_layout = False
    chart_line.legend.position = XL_LEGEND_POSITION.LEFT
    chart_line.legend.font.size = Pt(8)

    for line_serie in chart_line.series:
        line_serie.smooth = False
        line_serie.marker.style = XL_MARKER_STYLE.CIRCLE
        line_serie.data_labels.show_value = True
        line_serie.data_labels.number_format = '0.00%'
        line_serie.data_labels.font.size = Pt(8)

    value_axis_line = chart_line.value_axis
    value_axis_line.has_major_gridlines = False
    value_axis_line.major_tick_mark = XL_TICK_MARK.NONE
    value_axis_line.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis_line.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    value_axis_line.visible = False

    category_axis_line = chart_line.category_axis
    category_axis_line.has_major_gridlines = False
    category_axis_line.major_tick_mark = XL_TICK_MARK.NONE
    category_axis_line.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    category_axis_line.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    category_axis_line.visible = False

    # 设置堆积图y轴坐标和
    chart_data_cluster = CategoryChartData()
    chart_data_cluster.categories = data_years

    # 设置柱状堆积图的series数据--------------------------------
    series_Volumes1 = [vol / 1000 for vol in year_volumes_previous]
    series_Volumes2 = [vol / 1000 for vol in year_volumes_local]

    chart_data_cluster.add_series(PR_Status_previous, series_Volumes1)
    chart_data_cluster.add_series(PR_Status_local, series_Volumes2)

    x, y, cx, cy = Cm(1), Cm(5.5), Cm(24), Cm(5.5)
    graphic_frame_cluster = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_cluster
    )

    chart_cluster = graphic_frame_cluster.chart
    # chart_stack.has_title = True
    # chart_stack.chart_title.has_text_frame = True
    # chart_stack.chart_title.text_frame.text = "maoyadong"
    # chart_stack.chart_title.text_frame.paragraphs[0].font.size = Pt(10)

    chart_cluster.has_legend = True
    chart_cluster.legend.position = XL_LEGEND_POSITION.LEFT  # XL_LEGEND_POSITION.CORNER
    chart_cluster.legend.include_in_layout = False
    chart_cluster.legend.font.size = Pt(10)

    chart_cluster.series[0].data_labels.show_value = True
    chart_cluster.series[0].data_labels.number_format = '0'
    chart_cluster.series[0].data_labels.font.size = Pt(8)
    # chart_cluster.series[0].data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
    chart_cluster.series[1].data_labels.show_value = True
    chart_cluster.series[1].data_labels.number_format = '0'
    chart_cluster.series[1].data_labels.font.size = Pt(8)
    # chart_cluster.series[1].data_labels.position = XL_DATA_LABEL_POSITION.ABOVE


    value_axis_cluster = chart_cluster.value_axis
    value_axis_cluster.has_major_gridlines = False
    value_axis_cluster.major_tick_mark = XL_TICK_MARK.NONE
    value_axis_cluster.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis_cluster.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    value_axis_cluster.visible = False

    category_axis_cluster = chart_cluster.category_axis
    category_axis_cluster.has_major_gridlines = False
    category_axis_cluster.major_tick_mark = XL_TICK_MARK.NONE
    category_axis_cluster.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    category_axis_cluster.tick_labels.font.size = Pt(8)
    category_axis_cluster.format.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
    category_axis_cluster.visible = True

    # value_axis_stack.has_title = True
    # value_axis_stack.axis_title.has_text_frame = True
    # value_axis_stack.axis_title.text_frame.text = "False positive"
    # value_axis_stack.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

    # 开始创建表格
    rows = len(categories_order) + 2

    cols = len(data_years) + 1
    table_width = 24
    table_height = 2
    top = Cm(10.7)
    left = Cm(1.5)  # Inches(2.0)
    width = Cm(table_width)  # Inches(6.0)
    height = Cm(table_height)  # Inches(0.8)

    # 添加表格到幻灯片 --------------------
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # 给data_year加E
    data_years_E = [str(year) + 'E' if year > start_year else str(year) for year in data_years]

    # 设置单元格宽度
    columns_width = table_width / cols - 0.1
    for i in range(cols):
        if i == 0:
            table.columns[i].width = Cm(columns_width + 0.3)
        else:
            table.columns[i].width = Cm(columns_width)  # Inches(2.0)

    row_height = table_height / rows
    for i in range(rows):
        table.rows[i].height = Cm(row_height)  # Inches(2.0)

    # 设置标题行
    for i in range(cols):
        if i == 0:
            table.cell(0, i).text = "MS% Change"
        else:
            table.cell(0, i).text = str(data_years_E[i - 1])

    # 设置最后一行内容
    for i in range(cols):
        if i == 0:
            table.cell(rows - 1, i).text = "Group Total"
        else:
            table.cell(rows - 1, i).text = format(year_ms_change[i - 1], '.2%')

    # 填充表格第二行到倒数第二行数据
        for row_idx, category in enumerate(categories_order):
            table.cell(row_idx + 1, 0).text = category
            ms_changes = year_category_ms_change[category]
            for col_idx, change in enumerate(ms_changes):
                table.cell(row_idx + 1, col_idx + 1).text = format(change, '.2%')


    # 调整table每个cell的字体
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(table):
        if cell.text.strip() == '':
            cell.text = r'/'
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(7)

    # 开始添加注释文本框
    left = Cm(1)  # left，top为相对位置
    top = Cm(4)
    width = Cm(2)  # width，height为文本框的大小
    height = Cm(1)

    # 在指定位置添加文本框
    textbox = shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    # 在文本框中写入文字
    para = tf.add_paragraph()  # 新增段落
    para.text = "Volume\n'000units"  # 向段落写入文字
    para.line_spacing = 1.5  # 1.5 倍的行距
    para.font.size = Pt(8)

    prs.save('c:/auto-report/template_tmp7.pptx')
    print("maoyadong")
