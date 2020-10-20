import pandas as pd
from pandas import np
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE, XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Cm  # Inches
from pptx.util import Pt


class Report1(object):
    def __init__(self, data_file, config: ReportConfig):
        self.data_file, self.config = data_file, config

    def generate(self):
        # open ppt with cover
        df_vw = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', sheet_name=0)
        df_mkt = pd.read_excel(r'c:/auto-report/Database_small_demo.xlsx', sheet_name=0)
        # print(df_vw.head(10))
        print("df_vw=" + str(df_vw.shape[0]))
        # filter_data = df_vwGroup_baseData.query(' < @Cmean and B < @Cmean')
        # df_vwGroup_baseData.query('A < @Cmean and B < @Cmean')
        oem = df_vw.groupby('OEM').size().reset_index()["OEM"]
        brand = df_vw.groupby('Brand').size().reset_index()["Brand"]
        fuel_type = df_vw.groupby('Fuel_Type').size().reset_index()["Fuel_Type"]
        print(oem)
        print(brand)

        PR_Status = 'PR67.OP'
        start_year = 2018
        year_span = 8
        end_year = start_year + year_span
        # oem = ['FAW-VW', 'JAC-VW', 'JV TBD', 'SAIC-VW']
        # brand = ['Audi', 'Cupra', 'Jetta', 'Sihao', 'Skoda', 'VW']
        # 根据所有filter过滤大众数据
        df_vw_filter = df_vw[
            (df_vw['PR_Status'] == 'PR67.OP') & (df_vw['YEAR'] >= start_year) & (df_vw['YEAR'] <= end_year) & (
                df_vw['OEM'].isin(oem)) & (df_vw['Brand'].isin(brand))]
        print("df_vw_filter=" + str(df_vw_filter.shape[0]))

        # 全销售市场数据分组聚合-------------------------
        df_mkt_sum = df_mkt.groupby('Fuel_type').agg({'Volume': np.sum}).reset_index()

        # 大众销售市场数据分组聚合----------------------------
        df_vw_group_fuel = df_vw_filter.groupby(['Fuel_Type_Group']).size().reset_index()
        df_vw_group_fuel_year = df_vw_filter.groupby(['Fuel_Type_Group', 'YEAR']).agg({'Volume': np.sum}).reset_index()
        df_vw_group_year = df_vw_filter.groupby(['YEAR']).agg({'Volume': np.sum}).reset_index()
        df_vw_group_fuel_year.rename(columns={'Volume': 'total_Volume'}, inplace=True)
        df_vw_group_year.rename(columns={'Volume': 'total_Volume'}, inplace=True)
        print(df_vw_group_fuel_year)
        print(df_vw_group_year)
        print(df_vw_group_fuel)

        data_years = df_vw_group_year.sort_values(['YEAR'], ascending=[True])['YEAR']
        data_fuel_group = df_vw_group_fuel.sort_values(['Fuel_Type_Group'], ascending=[True])['Fuel_Type_Group']
        print(data_years)
        print(data_fuel_group)

        prs = Presentation('c:/auto-report/cover.pptx')
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        # 开始创建点线图-----------------------------------
        chart_data_line = ChartData()
        chart_data_line.categories = data_years
        mkt_volume = df_mkt_sum[(df_mkt_sum['Fuel_type'].isin(fuel_type))].agg({'Volume': np.sum}).reset_index().iloc[
            0, 1]
        total_volumes = df_vw_group_fuel_year[(df_vw_group_fuel_year['Fuel_Type_Group'].isin(data_fuel_group))] \
            .groupby(['YEAR']).agg({'total_Volume': np.sum}).reset_index() \
            .sort_values(['YEAR'], ascending=[True])['total_Volume']
        series_mkt_rate = [vol / mkt_volume for vol in total_volumes]
        chart_data_line.add_series('MS%', series_mkt_rate)

        x, y, cx, cy = Cm(1), Cm(5), Cm(24), Cm(4)
        chart_line = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_line
        ).chart

        chart_line.has_legend = True
        chart_line.legend.include_in_layout = False
        chart_line.legend.position = XL_LEGEND_POSITION.LEFT
        chart_line.legend.font.size = Pt(10)
        chart_line.chart_title.text_frame.clear()

        for line_serie in chart_line.series:
            line_serie.smooth = True
            line_serie.marker.style = XL_MARKER_STYLE.CIRCLE
            line_serie.data_labels.show_value = True
            line_serie.data_labels.number_format = '0.0%'
            line_serie.data_labels.font.size = Pt(10)
            line_serie.data_labels.position = XL_LABEL_POSITION.ABOVE

        value_axis_line = chart_line.value_axis
        value_axis_line.has_major_gridlines = False
        value_axis_line.major_tick_mark = XL_TICK_MARK.NONE
        value_axis_line.tick_label_position = XL_TICK_LABEL_POSITION.NONE
        value_axis_line.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        value_axis_line.visible = False

        category_axis_line = chart_line.category_axis
        category_axis_line.has_major_gridlines = False
        category_axis_line.major_tick_mark = XL_TICK_MARK.NONE
        category_axis_line.tick_label_position = XL_TICK_LABEL_POSITION.NONE
        category_axis_line.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        category_axis_line.visible = False

        # 开始创建柱状堆积图---------------------------------
        chart_data_stack = CategoryChartData()
        chart_data_stack.categories = data_years

        for fuel in data_fuel_group:
            series_Volumes = df_vw_group_fuel_year[(df_vw_group_fuel_year['Fuel_Type_Group'] == fuel)] \
                .sort_values(['YEAR'], ascending=[True])['total_Volume']
            series_Volumes = [vol / 1000 for vol in series_Volumes]
            chart_data_stack.add_series(fuel, series_Volumes)

        x, y, cx, cy = Cm(1), Cm(8), Cm(24), Cm(6)
        graphic_frame_stack = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data_stack
        )

        chart_stack = graphic_frame_stack.chart
        # chart_stack.has_title = True
        # chart_stack.chart_title.has_text_frame = True
        # chart_stack.chart_title.text_frame.text = "maoyadong"
        # chart_stack.chart_title.text_frame.paragraphs[0].font.size = Pt(10)

        chart_stack.has_legend = True
        chart_stack.legend.position = XL_LEGEND_POSITION.LEFT  # XL_LEGEND_POSITION.CORNER
        chart_stack.legend.include_in_layout = False
        chart_stack.legend.font.size = Pt(10)

        for line_serie in chart_stack.series:
            line_serie.smooth = True
            line_serie.data_labels.show_value = True
            line_serie.data_labels.number_format = '0'
            line_serie.data_labels.font.size = Pt(10)

        value_axis_stack = chart_stack.value_axis
        value_axis_stack.has_major_gridlines = False
        value_axis_stack.major_tick_mark = XL_TICK_MARK.NONE
        value_axis_stack.tick_label_position = XL_TICK_LABEL_POSITION.NONE
        value_axis_stack.format.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        value_axis_stack.visible = False

        category_axis_stack = chart_stack.category_axis
        category_axis_stack.has_major_gridlines = False
        category_axis_stack.major_tick_mark = XL_TICK_MARK.NONE
        category_axis_stack.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        category_axis_stack.tick_labels.font.size = Pt(10)
        category_axis_stack.format.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
        category_axis_stack.visible = True

        # value_axis_stack.has_title = True
        # value_axis_stack.axis_title.has_text_frame = True
        # value_axis_stack.axis_title.text_frame.text = "False positive"
        # value_axis_stack.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

        # 开始创建表格
        rows = 3
        cols = len(data_years) + 1
        table_width = 24
        table_height = 3
        top = Cm(14)
        left = Cm(1)  # Inches(2.0)
        width = Cm(table_width)  # Inches(6.0)
        height = Cm(table_height)  # Inches(0.8)

        # 添加表格到幻灯片 --------------------
        table = shapes.add_table(rows, cols, left, top, width, height).table

        # 给data_year加E
        data_years_E = [str(year) + 'E' if year > start_year else str(year) for year in data_years]

        # 设置单元格宽度
        columns_width = table_width / cols
        for i in range(cols):
            table.columns[i].width = Cm(columns_width)  # Inches(2.0)

        row_height = table_height / rows
        for i in range(rows):
            table.rows[i].height = Cm(row_height)  # Inches(2.0)

        # 设置标题行
        for i in range(cols):
            if i == 0:
                table.cell(0, i).text = 'MS%'
            else:
                table.cell(0, i).text = str(data_years_E[i - 1])

        # 设置行头

        # 填充表格数据
        for row_idx, fuel_group in enumerate(data_fuel_group):
            series_Volumes = df_vw_group_fuel_year[(df_vw_group_fuel_year['Fuel_Type_Group'] == fuel_group)] \
                .sort_values(['YEAR'], ascending=[True])['total_Volume']
            if fuel_group == 'ICE':
                mkt_volume = df_mkt_sum[(df_mkt_sum['Fuel_type'] == 'ICE')].reset_index().loc[0, 'Volume']
            else:
                mkt_volume = df_mkt_sum[(df_mkt_sum['Fuel_type'].isin(['BEV', 'PHEV']))].agg({'Volume': np.sum}) \
                    .reset_index().iloc[0, 1]
            print(mkt_volume)
            fuel_group_mkt_rate = [vol / mkt_volume for vol in series_Volumes]
            print(fuel_group_mkt_rate)
            table.cell(row_idx + 1, 0).text = fuel_group + ' MKT%'
            for col_idx, rate in enumerate(fuel_group_mkt_rate):
                table.cell(row_idx + 1, col_idx + 1).text = format(rate, '.1%')  # str(rate)
                # table.cell(row_idx + 1, col_idx + 1).text_frame.paragraphs[0].number_format = '0.0%'

        # 调整table每个cell的字体
        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)

        # 开始添加注释文本框
        left = Cm(1)  # left，top为相对位置
        top = Cm(4)
        width = Cm(2)  # width，height为文本框的大小
        height = Cm(1)

        # 在指定位置添加文本框
        textbox = shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        # 在文本框中写入文字
        para = tf.add_paragraph()  # 新增段落
        para.text = "Volume\n'000units"  # 向段落写入文字
        para.line_spacing = 1.5  # 1.5 倍的行距
        para.font.size = Pt(10)

        # 开始绘制普通柱状图 --------------------
        # chart_data = ChartData()
        # chart_data.categories = name_objects
        # chart_data.add_series(name_AIs[0], val_AI1)
        # chart_data.add_series(name_AIs[1], val_AI2)
        # chart_data.add_series(name_AIs[2], val_AI3)
        # chart_data.add_series(name_AIs[3], val_AI4)
        #
        #
        # x, y, cx, cy = Cm(0.5), Cm(6.2), Cm(24), Cm(6)
        # graphic_frame_cluster = slide.shapes.add_chart(
        #     XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        # )
        # chart_cluster = graphic_frame_cluster.chart
        # value_axis_cluster = chart_cluster.value_axis
        # category_axis_cluster = chart_cluster.category_axis
        # # value_axis.maximum_scale = 200.0
        # value_axis_cluster.has_major_gridlines = False
        # chart_cluster.value_axis.major_tick_mark = XL_TICK_MARK.NONE
        # category_axis_cluster.major_tick_mark = XL_TICK_MARK.NONE
        # # value_axis.has_minor_gridlines = False
        # value_axis_cluster.tick_label_position = XL_TICK_LABEL_POSITION.NONE
        # category_axis_cluster.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        #
        # category_axis_cluster.tick_labels.font.size = Pt(10)
        # category_axis_cluster.format.line.width = Cm(0)
        # value_axis_cluster.format.line.width = Cm(0)
        # value_axis_cluster.format.line.dash_style = MSO_LINE_DASH_STYLE.DASH_DOT
        # value_axis_cluster.visible = False
        # category_axis_cluster.visible = False
        return prs


if __name__ == '__main__':
    report1 = Report1(r'c:/auto-report/Database_small_demo.xlsx', None).generate().save(r'c:/auto-report/template_tmp.pptx')
